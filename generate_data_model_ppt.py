"""
AgentForge — ERwin Data Model PowerPoint Generator
ERD diagram drawn as native PPT shapes + connector lines (no external images)
Run: python generate_data_model_ppt.py
Output: AgentForge_DataModel.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colors ────────────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x11, 0x17)
BG_CARD      = RGBColor(0x1E, 0x21, 0x30)
BG_HEADER    = RGBColor(0x1A, 0x1D, 0x2E)
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_H     = RGBColor(0x4C, 0x1D, 0x95)
PURPLE_L     = RGBColor(0xC4, 0xB5, 0xFD)
BLUE         = RGBColor(0x3B, 0x82, 0xF6)
BLUE_H       = RGBColor(0x1E, 0x3A, 0x5F)
BLUE_L       = RGBColor(0x93, 0xC5, 0xFD)
TEAL         = RGBColor(0x05, 0x96, 0x69)
TEAL_H       = RGBColor(0x13, 0x4E, 0x4A)
TEAL_L       = RGBColor(0x5E, 0xEA, 0xD4)
ORANGE       = RGBColor(0xF9, 0x73, 0x16)
ORANGE_H     = RGBColor(0x43, 0x14, 0x07)
ORANGE_L     = RGBColor(0xFE, 0xD7, 0xAA)
GREEN        = RGBColor(0x16, 0xA3, 0x4A)
GREEN_H      = RGBColor(0x14, 0x53, 0x2D)
GREEN_L      = RGBColor(0x86, 0xEF, 0xAC)
YELLOW       = RGBColor(0xF5, 0x9E, 0x0B)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY         = RGBColor(0x94, 0xA3, 0xB8)
GRAY_DIM     = RGBColor(0x64, 0x74, 0x8B)
BORDER       = RGBColor(0x33, 0x41, 0x55)
ROW_ALT      = RGBColor(0x16, 0x18, 0x24)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Low-level helpers ─────────────────────────────────────────────────────────
def slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color=BG_CARD, line_color=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h, color=WHITE, size=11, bold=False,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.color.rgb = color; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    return txb

def add_connector(slide, x1, y1, x2, y2, color=BORDER, width=Pt(1.5), dash=True):
    """Draw a straight line connector."""
    cx = slide.shapes.add_connector(1, x1, y1, x2, y2)
    cx.line.color.rgb = color
    cx.line.width = width
    if dash:
        cx.line.dash_style = 4   # dash
    return cx

# ── ERwin entity card builder ─────────────────────────────────────────────────
ROW_H  = Inches(0.215)
HDR_H  = Inches(0.295)

def entity_height(fields):
    return HDR_H + len(fields) * ROW_H

def draw_entity(slide, title, fields, x, y, w,
                hdr_fill, hdr_text, border_color):
    """
    Draw a full ERwin-style entity box with header + rows.
    fields = list of (key, name, dtype)
      key: 'PK' | 'FK' | ''
    Returns (x, y, w, total_h) bounding box.
    """
    h = entity_height(fields)

    # outer border
    add_rect(slide, x, y, w, h, fill_color=BG_CARD,
             line_color=border_color, line_w=Pt(2))

    # header bar
    add_rect(slide, x, y, w, HDR_H, fill_color=hdr_fill)
    add_text(slide, title,
             x + Inches(0.1), y + Inches(0.04),
             w - Inches(0.12), HDR_H,
             color=hdr_text, size=8.5, bold=True)

    # separator line under header
    add_rect(slide, x, y + HDR_H, w, Pt(1.5), fill_color=border_color)

    # field rows
    for i, (key, fname, dtype) in enumerate(fields):
        ry = y + HDR_H + i * ROW_H
        row_bg = ROW_ALT if i % 2 == 0 else BG_CARD
        add_rect(slide, x, ry, w, ROW_H, fill_color=row_bg)

        # divider
        add_rect(slide, x, ry, w, Pt(0.5), fill_color=BORDER)

        # key badge
        key_color = YELLOW if key == 'PK' else (BLUE_L if key == 'FK' else GRAY_DIM)
        add_text(slide, key,
                 x + Inches(0.05), ry + Inches(0.02),
                 Inches(0.22), ROW_H,
                 color=key_color, size=6.5, bold=True)

        # field name
        add_text(slide, fname,
                 x + Inches(0.28), ry + Inches(0.02),
                 w * 0.52, ROW_H,
                 color=WHITE, size=7)

        # data type
        add_text(slide, dtype,
                 x + Inches(0.28) + w * 0.52, ry + Inches(0.02),
                 w * 0.43, ROW_H,
                 color=GRAY_DIM, size=6, italic=True)

    return x, y, w, h


def anchor(x, y, w, h, side='right', frac=0.5):
    """Return (px, py) connection point on entity edge."""
    if side == 'right':  return x + w, y + h * frac
    if side == 'left':   return x,     y + h * frac
    if side == 'bottom': return x + w * frac, y + h
    if side == 'top':    return x + w * frac, y


def draw_relation(slide, p1, p2, label='1:N', color=PURPLE_L):
    """Draw an elbow connector with a cardinality label."""
    x1, y1 = p1; x2, y2 = p2
    mx = (x1 + x2) / 2

    # two-segment elbow
    add_connector(slide, x1, y1, mx,  y1, color=color, width=Pt(1.5), dash=True)
    add_connector(slide, mx,  y1, mx,  y2, color=color, width=Pt(1.5), dash=True)
    add_connector(slide, mx,  y2, x2,  y2, color=color, width=Pt(1.5), dash=True)

    # arrowhead circle at target
    cr = Inches(0.04)
    add_rect(slide, x2 - cr, y2 - cr, cr*2, cr*2,
             fill_color=color, line_color=None)

    # label
    lx = mx + Inches(0.05)
    ly = min(y1, y2) + abs(y2 - y1) * 0.35
    add_text(slide, label, lx, ly, Inches(0.45), Inches(0.22),
             color=color, size=6.5, bold=True)


# =============================================================================
# SLIDE 1 — Cover
# =============================================================================
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)

add_rect(sl, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=PURPLE)
add_rect(sl, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), fill_color=PURPLE)

add_text(sl, 'AgentForge', Inches(1), Inches(1.4), Inches(11.3), Inches(1.3),
         color=PURPLE_L, size=52, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, 'ERwin Data Model  |  Enterprise AI Agent Builder Platform',
         Inches(1), Inches(2.65), Inches(11.3), Inches(0.5),
         color=GRAY, size=15, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(2.5), Inches(3.15), Inches(8.3), Inches(0.04), fill_color=BORDER)

add_text(sl, 'PostgreSQL  +  SQLAlchemy ORM  +  pgvector  +  Azure OpenAI GPT-4o',
         Inches(1), Inches(3.3), Inches(11.3), Inches(0.4),
         color=GRAY_DIM, size=11, align=PP_ALIGN.CENTER)

badges = [
    ('7 Entities',       PURPLE, PURPLE_H),
    ('6 FK Relations',   BLUE,   BLUE_H),
    ('3 User Roles',     TEAL,   TEAL_H),
    ('RAG Pipeline',     GREEN,  GREEN_H),
    ('Guardrails',       ORANGE, ORANGE_H),
]
bx = Inches(0.9)
for label, fc, bc in badges:
    bw = Inches(2.2)
    add_rect(sl, bx, Inches(3.9), bw, Inches(0.45),
             fill_color=bc, line_color=fc, line_w=Pt(1.5))
    add_text(sl, label, bx, Inches(3.97), bw, Inches(0.32),
             color=fc, size=10, bold=True, align=PP_ALIGN.CENTER)
    bx += bw + Inches(0.2)

add_text(sl, 'n.sureshmanikandan@accenture.com  |  July 2026',
         Inches(1), Inches(6.9), Inches(11.3), Inches(0.35),
         color=GRAY_DIM, size=9, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2 — Full ERwin Diagram (all 7 entities + relationships as shapes)
# =============================================================================
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, Inches(0), Inches(0), prs.slide_width, Inches(0.5), fill_color=BG_HEADER)
add_text(sl, 'ERwin Data Model — AgentForge Platform (All 7 Entities)',
         Inches(0.3), Inches(0.08), Inches(10), Inches(0.38),
         color=PURPLE_L, size=14, bold=True)
add_text(sl, 'Slide 2 of 7', Inches(11.8), Inches(0.08), Inches(1.5), Inches(0.35),
         color=GRAY_DIM, size=8, align=PP_ALIGN.RIGHT)

# ── Entity definitions ────────────────────────────────────────────────────────
users_f = [
    ('PK', 'id',              'VARCHAR(UUID)'),
    ('',   'email',           'VARCHAR UNIQUE'),
    ('',   'hashed_password', 'VARCHAR'),
    ('',   'full_name',       'VARCHAR'),
    ('',   'role',            'ENUM'),
    ('',   'is_active',       'BOOLEAN'),
    ('',   'created_at',      'DATETIME'),
]
agents_f = [
    ('PK', 'id',             'VARCHAR(UUID)'),
    ('',   'name',           'VARCHAR'),
    ('',   'description',    'VARCHAR'),
    ('',   'system_prompt',  'TEXT'),
    ('',   'model',          'VARCHAR'),
    ('',   'tools',          'JSON []'),
    ('',   'guardrails',     'JSON {}'),
    ('FK', 'created_by',     '-> users.id'),
    ('',   'current_version','INTEGER'),
    ('',   'created_at',     'DATETIME'),
    ('',   'updated_at',     'DATETIME'),
]
agentver_f = [
    ('PK', 'id',         'VARCHAR(UUID)'),
    ('FK', 'agent_id',   '-> agents.id'),
    ('',   'version',    'INTEGER'),
    ('',   'snapshot',   'JSON {}'),
    ('',   'created_at', 'DATETIME'),
]
workflow_f = [
    ('PK', 'id',          'VARCHAR(UUID)'),
    ('',   'name',        'VARCHAR'),
    ('',   'description', 'VARCHAR'),
    ('',   'nodes',       'JSON []'),
    ('',   'edges',       'JSON []'),
    ('FK', 'created_by',  '-> users.id'),
    ('',   'created_at',  'DATETIME'),
    ('',   'updated_at',  'DATETIME'),
]
kb_f = [
    ('PK', 'id',          'VARCHAR(UUID)'),
    ('',   'name',        'VARCHAR'),
    ('',   'description', 'VARCHAR'),
    ('FK', 'agent_id',    '-> agents.id'),
    ('FK', 'created_by',  '-> users.id'),
    ('',   'created_at',  'DATETIME'),
]
doc_f = [
    ('PK', 'id',          'VARCHAR(UUID)'),
    ('FK', 'kb_id',       '-> knowledge_bases.id'),
    ('',   'filename',    'VARCHAR'),
    ('',   'content',     'TEXT'),
    ('',   'chunk_count', 'INTEGER'),
    ('',   'status',      'VARCHAR'),
    ('',   'created_at',  'DATETIME'),
]
audit_f = [
    ('PK', 'id',                  'VARCHAR(UUID)'),
    ('FK', 'user_id',             '-> users.id'),
    ('',   'action',              'VARCHAR'),
    ('',   'resource_type',       'VARCHAR'),
    ('',   'resource_id',         'VARCHAR'),
    ('',   'input_snapshot',      'JSON {}'),
    ('',   'output_snapshot',     'JSON {}'),
    ('',   'guardrail_triggered', 'BOOLEAN'),
    ('',   'latency_ms',          'INTEGER'),
    ('',   'created_at',          'DATETIME'),
]

# ── Layout positions ──────────────────────────────────────────────────────────
EW = Inches(2.9)   # entity width

# Row 1
UX, UY   = Inches(0.18), Inches(0.58)   # users
AX, AY   = Inches(3.35), Inches(0.58)   # agents
AVAX, AVAY = Inches(6.55), Inches(0.58) # agent_versions

# Row 2
WX, WY   = Inches(0.18), Inches(3.90)   # workflows
KX, KY   = Inches(3.35), Inches(3.90)   # knowledge_bases
DX, DY   = Inches(6.55), Inches(3.90)   # documents
ALOX, ALOY = Inches(9.85), Inches(0.58) # audit_logs

# ── Draw all entities ─────────────────────────────────────────────────────────
_, _, _, UH  = draw_entity(sl, 'USERS',            users_f,    UX,   UY,   EW, PURPLE_H, PURPLE_L, PURPLE)
_, _, _, AH  = draw_entity(sl, 'AGENTS',           agents_f,   AX,   AY,   EW, BLUE_H,   BLUE_L,   BLUE)
_, _, _, AVH = draw_entity(sl, 'AGENT_VERSIONS',   agentver_f, AVAX, AVAY, EW, BLUE_H,   BLUE_L,   BLUE)
_, _, _, WH  = draw_entity(sl, 'WORKFLOWS',        workflow_f, WX,   WY,   EW, TEAL_H,   TEAL_L,   TEAL)
_, _, _, KH  = draw_entity(sl, 'KNOWLEDGE_BASES',  kb_f,       KX,   KY,   EW, GREEN_H,  GREEN_L,  GREEN)
_, _, _, DH  = draw_entity(sl, 'DOCUMENTS',        doc_f,      DX,   DY,   EW, GREEN_H,  GREEN_L,  GREEN)
_, _, _, ALH = draw_entity(sl, 'AUDIT_LOGS',       audit_f,    ALOX, ALOY, EW, ORANGE_H, ORANGE_L, ORANGE)

# ── Relationship connectors ───────────────────────────────────────────────────
# 1. users -> agents  (1:N, created_by)
draw_relation(sl,
    anchor(UX, UY, EW, UH, 'right', 0.55),
    anchor(AX, AY, EW, AH, 'left',  0.55),
    '1:N', PURPLE_L)

# 2. agents -> agent_versions  (1:N)
draw_relation(sl,
    anchor(AX, AY, EW, AH, 'right', 0.55),
    anchor(AVAX, AVAY, EW, AVH, 'left', 0.55),
    '1:N', BLUE_L)

# 3. users -> workflows  (1:N)
draw_relation(sl,
    anchor(UX, UY, EW, UH, 'bottom', 0.4),
    anchor(WX, WY, EW, WH, 'top',    0.4),
    '1:N', PURPLE_L)

# 4. users -> knowledge_bases  (1:N)
draw_relation(sl,
    anchor(UX, UY, EW, UH, 'bottom', 0.7),
    anchor(KX, KY, EW, KH, 'top',    0.3),
    '1:N', PURPLE_L)

# 5. agents -> knowledge_bases  (1:N, optional)
draw_relation(sl,
    anchor(AX, AY, EW, AH, 'bottom', 0.5),
    anchor(KX, KY, EW, KH, 'top',    0.6),
    '1:N', BLUE_L)

# 6. knowledge_bases -> documents  (1:N)
draw_relation(sl,
    anchor(KX, KY, EW, KH, 'right', 0.5),
    anchor(DX, DY, EW, DH, 'left',  0.5),
    '1:N', GREEN_L)

# 7. users -> audit_logs  (1:N)
draw_relation(sl,
    anchor(UX, UY, EW, UH, 'right', 0.25),
    anchor(ALOX, ALOY, EW, ALH, 'left', 0.25),
    '1:N', ORANGE_L)

# ── Legend ────────────────────────────────────────────────────────────────────
LX = Inches(9.85); LY = Inches(5.05)
add_rect(sl, LX, LY, Inches(3.3), Inches(2.1), fill_color=BG_CARD, line_color=BORDER, line_w=Pt(1))
add_text(sl, 'LEGEND', LX + Inches(0.15), LY + Inches(0.1),
         Inches(3.0), Inches(0.25), color=GRAY_DIM, size=8, bold=True)

legend_items = [
    ('PK',              YELLOW,   'Primary Key (UUID)'),
    ('FK',              BLUE_L,   'Foreign Key reference'),
    ('1:N',             GRAY,     'One-to-Many cardinality'),
    ('PURPLE border',   PURPLE_L, 'Auth / User domain'),
    ('BLUE border',     BLUE_L,   'Agent / Version domain'),
    ('TEAL/GREEN',      TEAL_L,   'Workflow / RAG domain'),
    ('ORANGE border',   ORANGE_L, 'Audit / Observability'),
]
for i, (badge, c, desc) in enumerate(legend_items):
    ly2 = LY + Inches(0.42) + i * Inches(0.22)
    add_rect(sl, LX + Inches(0.15), ly2, Inches(0.04), Inches(0.14), fill_color=c)
    add_text(sl, desc, LX + Inches(0.28), ly2 - Inches(0.02),
             Inches(2.9), Inches(0.22), color=GRAY, size=7)


# =============================================================================
# SLIDE 3 — ERD Detail: Auth & Agent Core (zoomed entity cards + notes)
# =============================================================================
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, Inches(0), Inches(0), prs.slide_width, Inches(0.5), fill_color=BG_HEADER)
add_text(sl, 'ERD Detail — Auth & Agent Core (users, agents, agent_versions)',
         Inches(0.3), Inches(0.08), Inches(10), Inches(0.38),
         color=PURPLE_L, size=14, bold=True)
add_text(sl, 'Slide 3 of 7', Inches(11.8), Inches(0.08), Inches(1.5), Inches(0.35),
         color=GRAY_DIM, size=8, align=PP_ALIGN.RIGHT)

EW3 = Inches(3.8)
GX = Inches(0.3)
_, _, _, UH3  = draw_entity(sl, 'USERS',          users_f,    GX,             Inches(0.6), EW3, PURPLE_H, PURPLE_L, PURPLE)
_, _, _, AH3  = draw_entity(sl, 'AGENTS',         agents_f,   GX + Inches(4.3), Inches(0.6), EW3, BLUE_H,   BLUE_L,   BLUE)
_, _, _, AVH3 = draw_entity(sl, 'AGENT_VERSIONS', agentver_f, GX + Inches(8.6), Inches(0.6), EW3, BLUE_H,   BLUE_L,   BLUE)

# connectors
draw_relation(sl,
    anchor(GX, Inches(0.6), EW3, UH3, 'right', 0.55),
    anchor(GX + Inches(4.3), Inches(0.6), EW3, AH3, 'left', 0.55),
    '1:N', PURPLE_L)
draw_relation(sl,
    anchor(GX + Inches(4.3), Inches(0.6), EW3, AH3, 'right', 0.55),
    anchor(GX + Inches(8.6), Inches(0.6), EW3, AVH3, 'left', 0.55),
    '1:N', BLUE_L)

# Notes
add_rect(sl, Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.65), fill_color=BG_CARD, line_color=BORDER)
add_text(sl, 'Entity Notes', Inches(0.5), Inches(5.6), Inches(5), Inches(0.3),
         color=GRAY, size=9, bold=True)
for i, note in enumerate([
    '  users.role is a PostgreSQL ENUM — values: admin | developer | viewer. Enforced at JWT decode and API middleware.',
    '  agents.tools is a JSON array of tool name strings (email, slack, github, web_search, calculator) — extensible without schema migrations.',
    '  agents.guardrails is JSON: {"pii": true, "hallucination": true, "max_tokens": 2048} — drives GuardrailsEngine at runtime.',
    '  agent_versions.snapshot stores a frozen JSON copy of the full agent config — enables exact point-in-time rollback.',
    '  All PKs are UUID (uuid4) generated in Python — portable across environments, no DB sequence dependency.',
]):
    add_text(sl, note, Inches(0.5), Inches(5.9) + i * Inches(0.24),
             Inches(12.5), Inches(0.24), color=GRAY, size=8)


# =============================================================================
# SLIDE 4 — ERD Detail: Workflows & RAG Pipeline
# =============================================================================
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, Inches(0), Inches(0), prs.slide_width, Inches(0.5), fill_color=BG_HEADER)
add_text(sl, 'ERD Detail — Workflows & RAG Pipeline (workflows, knowledge_bases, documents)',
         Inches(0.3), Inches(0.08), Inches(10.5), Inches(0.38),
         color=TEAL_L, size=14, bold=True)
add_text(sl, 'Slide 4 of 7', Inches(11.8), Inches(0.08), Inches(1.5), Inches(0.35),
         color=GRAY_DIM, size=8, align=PP_ALIGN.RIGHT)

EW4 = Inches(3.8)
_, _, _, WH4 = draw_entity(sl, 'WORKFLOWS',       workflow_f, Inches(0.3),  Inches(0.6), EW4, TEAL_H,  TEAL_L,  TEAL)
_, _, _, KH4 = draw_entity(sl, 'KNOWLEDGE_BASES', kb_f,       Inches(4.6),  Inches(0.6), EW4, GREEN_H, GREEN_L, GREEN)
_, _, _, DH4 = draw_entity(sl, 'DOCUMENTS',       doc_f,      Inches(8.9),  Inches(0.6), EW4, GREEN_H, GREEN_L, GREEN)

draw_relation(sl,
    anchor(Inches(4.6), Inches(0.6), EW4, KH4, 'right', 0.5),
    anchor(Inches(8.9), Inches(0.6), EW4, DH4, 'left',  0.5),
    '1:N', GREEN_L)

# RAG pipeline flow diagram
add_text(sl, 'RAG Ingestion & Query Pipeline', Inches(0.3), Inches(4.2),
         Inches(8), Inches(0.3), color=GREEN_L, size=10, bold=True)

pipeline = [
    ('Upload\nDocument',     GREEN),
    ('LangChain\nSplitter',  TEAL),
    ('Embeddings\n(GPT)',    BLUE),
    ('Azure AI\nSearch / pgvector', PURPLE),
    ('Semantic\nRetrieval',  TEAL),
    ('GPT-4o\nAnswer',       GREEN),
]
for i, (step, c) in enumerate(pipeline):
    px = Inches(0.3) + i * Inches(2.12)
    add_rect(sl, px, Inches(4.6), Inches(1.92), Inches(0.72),
             fill_color=BG_CARD, line_color=c, line_w=Pt(1.5))
    add_text(sl, step, px + Inches(0.06), Inches(4.66),
             Inches(1.8), Inches(0.6), color=c, size=8, bold=True, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text(sl, '>>', px + Inches(1.92), Inches(4.82),
                 Inches(0.18), Inches(0.28), color=GRAY_DIM, size=9, bold=True)

add_rect(sl, Inches(0.3), Inches(5.55), Inches(12.7), Inches(1.6), fill_color=BG_CARD, line_color=BORDER)
add_text(sl, 'RAG & Workflow Notes', Inches(0.5), Inches(5.65), Inches(5), Inches(0.28),
         color=GRAY, size=9, bold=True)
for i, note in enumerate([
    '  workflows.nodes & edges store ReactFlow JSON arrays — the full visual canvas is serialised into a single PostgreSQL column.',
    '  knowledge_bases.agent_id is NULLABLE — a KB can be standalone (shared) or bound to a single agent.',
    '  documents.status lifecycle: processing (ingestion started) -> ready (chunks indexed) -> failed (error).',
    '  Document chunks are embedded and written to Azure AI Search index; pgvector is the fallback for local dev.',
    '  RAGEngine._retrieve() performs keyword scoring locally; in production replace with AzureSearchClient vector query.',
]):
    add_text(sl, note, Inches(0.5), Inches(5.92) + i * Inches(0.24),
             Inches(12.5), Inches(0.24), color=GRAY, size=8)


# =============================================================================
# SLIDE 5 — ERD Detail: Audit Logs + Full FK Reference Table
# =============================================================================
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, Inches(0), Inches(0), prs.slide_width, Inches(0.5), fill_color=BG_HEADER)
add_text(sl, 'ERD Detail — Audit Logs & Complete FK Reference',
         Inches(0.3), Inches(0.08), Inches(10), Inches(0.38),
         color=ORANGE_L, size=14, bold=True)
add_text(sl, 'Slide 5 of 7', Inches(11.8), Inches(0.08), Inches(1.5), Inches(0.35),
         color=GRAY_DIM, size=8, align=PP_ALIGN.RIGHT)

_, _, _, ALH5 = draw_entity(sl, 'AUDIT_LOGS', audit_f,
                             Inches(0.3), Inches(0.6), Inches(5.0),
                             ORANGE_H, ORANGE_L, ORANGE)

# FK reference table (right side)
add_rect(sl, Inches(5.7), Inches(0.6), Inches(7.3), Inches(3.8),
         fill_color=BG_CARD, line_color=BORDER)
add_rect(sl, Inches(5.7), Inches(0.6), Inches(7.3), Inches(0.32), fill_color=BG_HEADER)
add_text(sl, 'Complete Foreign-Key Reference', Inches(5.85), Inches(0.64),
         Inches(7.0), Inches(0.28), color=GRAY, size=8.5, bold=True)

fk_rows = [
    ('agents.created_by',          'users.id',           '1:N', PURPLE_L),
    ('workflows.created_by',       'users.id',           '1:N', PURPLE_L),
    ('knowledge_bases.created_by', 'users.id',           '1:N', PURPLE_L),
    ('audit_logs.user_id',         'users.id',           '1:N', ORANGE_L),
    ('agent_versions.agent_id',    'agents.id',          '1:N', BLUE_L),
    ('knowledge_bases.agent_id',   'agents.id (NULL OK)','1:N', BLUE_L),
    ('documents.kb_id',            'knowledge_bases.id', '1:N', GREEN_L),
]
hdr_cols = ['FK Column', 'References', 'Card']
hxs = [Inches(5.85), Inches(8.55), Inches(11.1)]
hws = [Inches(2.6),  Inches(2.45), Inches(0.65)]

for h, hx, hw in zip(hdr_cols, hxs, hws):
    add_text(sl, h, hx, Inches(1.0), hw, Inches(0.25),
             color=GRAY_DIM, size=7.5, bold=True)

for i, (fk, ref, card, c) in enumerate(fk_rows):
    ry = Inches(1.28) + i * Inches(0.36)
    bg = ROW_ALT if i % 2 == 0 else BG_CARD
    add_rect(sl, Inches(5.7), ry, Inches(7.3), Inches(0.36), fill_color=bg)
    vals = [fk, ref, card]
    clrs = [c, WHITE, GRAY]
    for v, hx, hw, vc in zip(vals, hxs, hws, clrs):
        add_text(sl, v, hx, ry + Inches(0.04), hw, Inches(0.3),
                 color=vc, size=7.5, bold=(vc == c))

# Notes
add_rect(sl, Inches(0.3), Inches(5.2), Inches(12.7), Inches(1.95), fill_color=BG_CARD, line_color=BORDER)
add_text(sl, 'Audit Log Design Notes', Inches(0.5), Inches(5.3), Inches(5), Inches(0.28),
         color=GRAY, size=9, bold=True)
for i, note in enumerate([
    '  audit_logs.user_id is NULLABLE — system/background actions (e.g. scheduled tasks) log with NULL user_id.',
    '  Every agents/{id}/run call appends one audit_logs row — full traceability: input, output, guardrail, latency.',
    '  guardrail_triggered = TRUE flags rows where Presidio detected PII or hallucination phrases were found.',
    '  resource_type + resource_id together form a polymorphic reference ("agent","wf-123") — no ENUM constraint needed.',
    '  The /control-plane/stats endpoint aggregates this table live: total_runs, avg_latency_ms, guardrail_count.',
    '  Recommended index: (resource_type, resource_id, created_at DESC) for per-agent log queries.',
]):
    add_text(sl, note, Inches(0.5), Inches(5.6) + i * Inches(0.24),
             Inches(12.5), Inches(0.24), color=GRAY, size=8)


# =============================================================================
# SLIDE 6 — Data Dictionary
# =============================================================================
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, Inches(0), Inches(0), prs.slide_width, Inches(0.5), fill_color=BG_HEADER)
add_text(sl, 'Data Dictionary — Column Reference (Key Columns)',
         Inches(0.3), Inches(0.08), Inches(10), Inches(0.38),
         color=WHITE, size=14, bold=True)
add_text(sl, 'Slide 6 of 7', Inches(11.8), Inches(0.08), Inches(1.5), Inches(0.35),
         color=GRAY_DIM, size=8, align=PP_ALIGN.RIGHT)

dict_entries = [
    ('users',           'role',               'ENUM',    'admin | developer | viewer  — enforced at JWT decode & FastAPI Depends', PURPLE_L),
    ('agents',          'tools',              'JSON []', 'List of tool names: [email, slack, github, web_search, calculator]',     BLUE_L),
    ('agents',          'guardrails',         'JSON {}', '{"pii": bool, "hallucination": bool, "max_tokens": int}',                BLUE_L),
    ('agent_versions',  'snapshot',           'JSON {}', 'Full frozen agent config at save time — enables rollback to any version',BLUE_L),
    ('workflows',       'nodes',              'JSON []', 'ReactFlow node objects: {id, type, position, data} per canvas node',     TEAL_L),
    ('workflows',       'edges',              'JSON []', 'ReactFlow edge objects: {id, source, target} — directed graph edges',    TEAL_L),
    ('knowledge_bases', 'agent_id',           'VARCHAR', 'NULLABLE FK — NULL = standalone KB shared across agents',                GREEN_L),
    ('documents',       'status',             'VARCHAR', 'processing | ready | failed — drives upload progress UI indicator',      GREEN_L),
    ('documents',       'content',            'TEXT',    'Raw extracted text; chunks indexed in Azure AI Search / pgvector',       GREEN_L),
    ('audit_logs',      'action',             'VARCHAR', "Dot-notation: 'agent.run', 'agent.update', 'kb.upload'",                 ORANGE_L),
    ('audit_logs',      'input_snapshot',     'JSON {}', 'Full request payload at call time — persisted for compliance & replay', ORANGE_L),
    ('audit_logs',      'guardrail_triggered','BOOLEAN', 'TRUE if PII redacted OR hallucination phrase detected in response',      ORANGE_L),
]

add_rect(sl, Inches(0.3), Inches(0.58), Inches(12.7), Inches(0.3), fill_color=BG_HEADER)
for h, hx, hw in zip(['Entity','Column','Type','Description'],
                      [Inches(0.4), Inches(2.05), Inches(3.45), Inches(4.55)],
                      [Inches(1.55), Inches(1.3), Inches(1.0), Inches(8.5)]):
    add_text(sl, h, hx, Inches(0.6), hw, Inches(0.26), color=GRAY_DIM, size=8, bold=True)

for i, (entity, col, dtype, desc, c) in enumerate(dict_entries):
    ry = Inches(0.9) + i * Inches(0.5)
    add_rect(sl, Inches(0.3), ry, Inches(12.7), Inches(0.48),
             fill_color=ROW_ALT if i % 2 == 0 else BG_CARD)
    for v, hx, hw, vc, vb in zip(
        [entity, col, dtype, desc],
        [Inches(0.4), Inches(2.05), Inches(3.45), Inches(4.55)],
        [Inches(1.55), Inches(1.3), Inches(1.0), Inches(8.5)],
        [c, WHITE, GRAY_DIM, GRAY],
        [True, True, False, False]
    ):
        add_text(sl, v, hx, ry + Inches(0.06), hw, Inches(0.38),
                 color=vc, size=7.5, bold=vb)


# =============================================================================
# SLIDE 7 — Summary & Roadmap
# =============================================================================
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, Inches(0), Inches(0), prs.slide_width, Inches(0.5), fill_color=BG_HEADER)
add_text(sl, 'Summary — All Entities & Implementation Roadmap',
         Inches(0.3), Inches(0.08), Inches(10), Inches(0.38),
         color=WHITE, size=14, bold=True)
add_text(sl, 'Slide 7 of 7', Inches(11.8), Inches(0.08), Inches(1.5), Inches(0.35),
         color=GRAY_DIM, size=8, align=PP_ALIGN.RIGHT)

entities_summary = [
    ('USERS',           '7 cols',  'Auth / RBAC — 3 roles (admin, developer, viewer)',          PURPLE, PURPLE_L),
    ('AGENTS',          '11 cols', 'Core entity — GPT-4o config, tools, guardrails, versioning',BLUE,   BLUE_L),
    ('AGENT_VERSIONS',  '5 cols',  'Immutable snapshots — full config rollback capability',     BLUE,   BLUE_L),
    ('WORKFLOWS',       '8 cols',  'ReactFlow canvas state — nodes + edges stored as JSON',     TEAL,   TEAL_L),
    ('KNOWLEDGE_BASES', '6 cols',  'RAG container — optional agent binding, standalone OK',     GREEN,  GREEN_L),
    ('DOCUMENTS',       '7 cols',  'Uploaded files — tracking status, content, chunk count',    GREEN,  GREEN_L),
    ('AUDIT_LOGS',      '10 cols', 'Full observability — every run, latency, guardrail result', ORANGE, ORANGE_L),
]
add_text(sl, 'All 7 Entities', Inches(0.3), Inches(0.62),
         Inches(6), Inches(0.28), color=GRAY, size=9, bold=True)
for i, (name, cols, desc, hc, tc) in enumerate(entities_summary):
    ey = Inches(0.95) + i * Inches(0.74)
    add_rect(sl, Inches(0.3), ey, Inches(6.1), Inches(0.66), fill_color=BG_CARD, line_color=hc, line_w=Pt(1.5))
    add_rect(sl, Inches(0.3), ey, Inches(0.1), Inches(0.66), fill_color=hc)
    add_text(sl, name, Inches(0.5), ey + Inches(0.06), Inches(2.5), Inches(0.28), color=tc, size=8.5, bold=True)
    add_text(sl, cols, Inches(3.1), ey + Inches(0.06), Inches(0.9), Inches(0.28), color=GRAY_DIM, size=7.5)
    add_text(sl, desc, Inches(0.5), ey + Inches(0.34), Inches(5.8), Inches(0.28), color=GRAY, size=7.5)

phases = [
    ('Phase 1',   'Project Bootstrap',            'FastAPI + Docker + React scaffold',         TEAL),
    ('Phase 2',   'DB Models + Auth + RBAC',       'SQLAlchemy ORM + JWT + Alembic migrations', BLUE),
    ('Phase 3',   'Azure OpenAI + NL Generator',   'GPT-4o wrapper + Prompt-to-Agent',          PURPLE),
    ('Phase 4',   'Guardrails Engine',             'Presidio PII + hallucination detection',    ORANGE),
    ('Phase 5',   'Multi-Agent Orchestrator',      'Manager/worker pattern + streaming',        GREEN),
    ('Phase 6-8', 'Agents API + RAG + Simulation', 'CRUD, KB upload, batch test runner',        TEAL),
    ('Phase 9-10','Tool Registry + Control Plane', 'Tool definitions + audit logs + stats',     BLUE),
    ('Phase 11-12','React UI + Docker + Wiring',   'Canvas, Dashboard, full stack deploy',      PURPLE),
]
add_text(sl, 'Implementation Roadmap', Inches(6.8), Inches(0.62),
         Inches(6), Inches(0.28), color=GRAY, size=9, bold=True)
for i, (ph, title, detail, c) in enumerate(phases):
    py = Inches(0.95) + i * Inches(0.68)
    add_rect(sl, Inches(6.8), py, Inches(6.2), Inches(0.6), fill_color=BG_CARD, line_color=c, line_w=Pt(1.2))
    add_text(sl, ph, Inches(6.92), py + Inches(0.05), Inches(0.9), Inches(0.25), color=c, size=7.5, bold=True)
    add_text(sl, title, Inches(7.85), py + Inches(0.05), Inches(4.9), Inches(0.25), color=WHITE, size=8, bold=True)
    add_text(sl, detail, Inches(6.92), py + Inches(0.3), Inches(5.9), Inches(0.25), color=GRAY, size=7.5)

add_rect(sl, Inches(0), Inches(7.18), prs.slide_width, Inches(0.32), fill_color=BG_HEADER)
add_text(sl, 'AgentForge  |  ERwin Data Model  |  n.sureshmanikandan@accenture.com  |  July 2026',
         Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.28),
         color=GRAY_DIM, size=8, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = r'C:\Users\n.sureshmanikandan\Repo1\AgentForge\AgentForge_DataModel.pptx'
prs.save(out)
print(f'DONE: {out}')
