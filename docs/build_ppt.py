"""
AgentForge Presentation Builder — V7.0
Run: python docs/build_ppt.py
Output: docs/AgentForge-PresentationV7.0.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT     = RGBColor(0x60, 0x6E, 0xF5)   # indigo/violet
ACCENT2    = RGBColor(0x10, 0xB9, 0x81)   # emerald
ACCENT3    = RGBColor(0xF5, 0x9E, 0x0B)   # amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
CARD_BG    = RGBColor(0x1E, 0x2A, 0x45)
DARK_TEXT  = RGBColor(0x1E, 0x29, 0x3B)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helper functions ──────────────────────────────────────────────────────────

def bg(slide, color=DARK_BG):
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def rect(slide, x, y, w, h, fill, radius=False):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def accent_bar(slide, y=Inches(0.08), color=ACCENT):
    r = rect(slide, 0, y, W, Inches(0.05), color)
    return r

def slide_num(slide, n):
    txt(slide, str(n), W - Inches(0.6), H - Inches(0.4), Inches(0.5), Inches(0.35),
        size=11, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def chip(slide, label, x, y, w=Inches(1.8), h=Inches(0.42), fill=ACCENT, tsize=13):
    rect(slide, x, y, w, h, fill)
    txt(slide, label, x, y, w, h, size=tsize, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def section_header(slide, title, subtitle=""):
    bg(slide)
    accent_bar(slide)
    rect(slide, Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.06), ACCENT)
    txt(slide, title,    Inches(0.5), Inches(2.2), Inches(12), Inches(1.2), size=44, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        txt(slide, subtitle, Inches(0.5), Inches(3.5), Inches(12), Inches(0.8), size=20,
            color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, title, body_lines, title_color=ACCENT, icon=""):
    rect(slide, x, y, w, h, CARD_BG)
    title_str = f"{icon}  {title}" if icon else title
    txt(slide, title_str, x + Inches(0.18), y + Inches(0.15), w - Inches(0.3), Inches(0.5),
        size=15, bold=True, color=title_color)
    body = "\n".join(body_lines)
    txt(slide, body, x + Inches(0.18), y + Inches(0.58), w - Inches(0.3), h - Inches(0.7),
        size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, Inches(0), ACCENT)

# Gradient overlay strip
rect(sl, 0, Inches(1.5), W, Inches(4.5), RGBColor(0x1A, 0x24, 0x3E))

# Big logo pill
chip(sl, "⚡  AgentForge", Inches(0.5), Inches(1.8), Inches(3.2), Inches(0.55), ACCENT, 16)

txt(sl, "Enterprise AI Agent Platform",
    Inches(0.4), Inches(2.6), Inches(9), Inches(1.4),
    size=48, bold=True, color=WHITE)

txt(sl, "Build · Orchestrate · Deploy · Govern — AI Agents at Scale",
    Inches(0.4), Inches(4.1), Inches(10), Inches(0.7),
    size=22, color=LIGHT_GRAY)

# Tags row
for i, tag in enumerate(["Azure OpenAI", "FastAPI", "React 18", "PostgreSQL + pgvector", "OpenTelemetry"]):
    chip(sl, tag, Inches(0.4 + i * 2.5), Inches(5.0), Inches(2.3), Inches(0.38),
         CARD_BG, 12)

txt(sl, "Version 7.0  |  July 2026  |  Confidential",
    Inches(0.4), Inches(6.2), Inches(6), Inches(0.4), size=12, color=LIGHT_GRAY)
slide_num(sl, 1)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Executive Summary", Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
    size=28, bold=True, color=WHITE)

summary = (
    "AgentForge is an enterprise-grade AI Agent Platform that empowers Accenture developers and clients to "
    "design, build, simulate, deploy, and govern intelligent AI agents at scale — without deep ML expertise.\n\n"
    "Powered by Azure OpenAI GPT-5.4-mini, it delivers a full-stack solution: a conversational Planning Architect "
    "(Prompt-to-Agent), a visual Workflow Builder, a built-in RAG pipeline, enterprise guardrails (PII + "
    "hallucination detection), multi-agent orchestration, voice agent support, and a real-time Control Plane — "
    "all with RBAC, audit logs, and OpenTelemetry observability baked in."
)
txt(sl, summary, Inches(0.5), Inches(1.1), Inches(12.4), Inches(2.2),
    size=14, color=LIGHT_GRAY)

# KPI row
kpis = [
    ("16 API Modules", "End-to-end coverage"),
    ("24 Frontend Pages", "Full SPA experience"),
    ("10 DB Tables", "Relational + vector store"),
    ("Zero ML Setup", "Prompt-to-agent in <2 min"),
    ("Enterprise Ready", "RBAC · Audit · OTEL"),
]
card_w = Inches(2.4)
for i, (k, v) in enumerate(kpis):
    x = Inches(0.4 + i * 2.58)
    rect(sl, x, Inches(3.55), card_w, Inches(1.6), CARD_BG)
    txt(sl, k, x + Inches(0.12), Inches(3.7), card_w - Inches(0.2), Inches(0.65),
        size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(sl, v, x + Inches(0.12), Inches(4.35), card_w - Inches(0.2), Inches(0.6),
        size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(sl, "Designed for Accenture | Delivered with Azure | Governed by Design",
    Inches(0.5), Inches(5.4), Inches(12), Inches(0.5),
    size=14, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
slide_num(sl, 2)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "The Problem We Solve", Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
    size=28, bold=True, color=WHITE)

problems = [
    ("🧩", "Fragmented Tooling",
     "Teams use disparate tools — LangChain, custom scripts, ad-hoc prompts — with no unified governance or reuse."),
    ("⏱️", "Slow Time-to-Agent",
     "Building, testing, and deploying even a simple RAG agent takes days of boilerplate coding and infra setup."),
    ("🔓", "No Enterprise Controls",
     "No PII redaction, no hallucination checks, no audit trails — violating compliance and data privacy requirements."),
    ("📉", "Zero Observability",
     "Once deployed, agents are black boxes. No tracing, no latency metrics, no ability to tune post-deployment."),
    ("🔁", "No Reusability",
     "Prompts, tools, and agent configs are siloed per project — no marketplace or shared blueprint library."),
    ("🗣️", "No Voice/Multi-modal",
     "Enterprises need voice-capable agents and multi-agent pipelines; existing DIY approaches lack this natively."),
]
for i, (icon, title, body) in enumerate(problems):
    row, col = divmod(i, 3)
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.1 + row * 2.0)
    card(sl, x, y, Inches(4.0), Inches(1.75), title, [body], ACCENT3, icon)
slide_num(sl, 3)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — PLATFORM OVERVIEW (Architecture Visual)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Platform Architecture Overview", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
    size=28, bold=True, color=WHITE)

layers = [
    (ACCENT,  "FRONTEND  (React 18 + Vite + TypeScript + TailwindCSS)",
     "Planning Architect · Agent Studio · Workflow Builder · RAG/KB · Playground · Voice · Safety · Evaluations · Control Plane"),
    (RGBColor(0x7C,0x3A,0xED), "API GATEWAY  (FastAPI 0.115 · JWT/RBAC · OpenTelemetry)",
     "Auth · Agents · RAG · Tools · Builder · Simulation · Control Plane · Safety · Evaluations · Voice · Architect · Team · API Keys"),
    (RGBColor(0x0E,0x79,0xB2), "CORE ENGINE  (Python 3.12)",
     "Prompt-to-Agent · Multi-Agent Orchestrator · RAG Engine (FAISS + pgvector) · Guardrails (Presidio) · Simulation · Tool Registry · Telemetry"),
    (ACCENT2, "AI SERVICES  (Azure OpenAI)",
     "GPT-5.4-mini (gpt-5.4-mini) · GPT-5.4-mini (gpt-5.4-mini) · Embeddings · Azure AI Search · Azure Speech (STT/TTS)"),
    (RGBColor(0xDC,0x26,0x26), "DATA LAYER  (PostgreSQL 16 + pgvector + SQLAlchemy 2.0 async)",
     "users · agents · agent_versions · workflows · workflow_runs · knowledge_bases · documents · audit_logs · api_keys · voice_sessions"),
]
for i, (clr, title, body) in enumerate(layers):
    y = Inches(1.1 + i * 1.15)
    rect(sl, Inches(0.4), y, Inches(12.4), Inches(1.0), CARD_BG)
    rect(sl, Inches(0.4), y, Inches(0.12), Inches(1.0), clr)
    txt(sl, title, Inches(0.65), y + Inches(0.05), Inches(12), Inches(0.42),
        size=13, bold=True, color=clr)
    txt(sl, body, Inches(0.65), y + Inches(0.45), Inches(12), Inches(0.42),
        size=11, color=LIGHT_GRAY)
slide_num(sl, 4)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — KEY FEATURES (12 features)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Core Feature Set", Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
    size=28, bold=True, color=WHITE)

features = [
    ("🧠", "Planning Architect (Prompt-to-Agent)",
     "Conversational AI that converts a plain-English description into a complete agent stack — plan, agents, tools, DB schema."),
    ("🔧", "Visual Workflow Builder",
     "Drag-and-drop canvas (React Flow) to wire manager + worker agents, tools, and conditions into pipelines. Full execution trace."),
    ("📚", "RAG Pipeline & Knowledge Base",
     "Upload PDF/DOCX/TXT → chunked → FAISS + pgvector → semantic retrieval. Per-agent KB scoping with Azure AI Search fallback."),
    ("🛡️", "Enterprise Guardrails",
     "PII redaction (Presidio — SSN, email, CC, phone) + hallucination detection on every agent response. Per-agent toggle."),
    ("🤝", "Multi-Agent Orchestration",
     "Manager agent fans out to N worker agents concurrently. Aggregates results. Supports sequential & parallel patterns."),
    ("🗣️", "Voice Agents",
     "Azure Speech STT → GPT-5.4-mini → Azure TTS pipeline. REST API for voice session management. Real-time audio I/O."),
    ("🧪", "Simulation & Evaluation Engine",
     "Run test scenarios against agents before deployment. Score output quality. Compare versions. Regression test suite."),
    ("🔒", "RBAC & API Key Management",
     "Three roles: Admin · Developer · Viewer. Per-user JWT auth. Scoped API keys for programmatic agent invocation."),
    ("📊", "Control Plane & Observability",
     "Live workflow status, modify running agents, pause/resume. OpenTelemetry traces + Azure Monitor + custom dashboards."),
    ("📝", "Audit Logs & Traceability",
     "Every agent action logged: user, action, resource, input/output snapshot, guardrail trigger, latency, trace ID."),
    ("🗂️", "Prompt Library & Blueprints",
     "Curated prompt templates and reusable agent blueprints for common use cases (support, analytics, RAG Q&A)."),
    ("⚡", "Architect Download Engine",
     "One-click export of deployable RAG Scaffold or Custom Code ZIP (React + FastAPI) — production-ready in minutes."),
]
for i, (icon, title, body) in enumerate(features):
    row, col = divmod(i, 4)
    x = Inches(0.32 + col * 3.26)
    y = Inches(1.05 + row * 1.95)
    card(sl, x, y, Inches(3.1), Inches(1.82), title, [body], ACCENT, icon)
slide_num(sl, 5)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — ROI & BUSINESS VALUE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "ROI & Business Value", Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
    size=28, bold=True, color=WHITE)

rois = [
    ("⚡ Time-to-Agent",    "From weeks → <2 min",  "85%↓",  ACCENT,
     "Prompt-to-Agent converts NL description into full agent stack instantly — no boilerplate"),
    ("💰 Dev Cost Savings", "~40 hrs/agent saved",   "$18K+", ACCENT2,
     "Reusable blueprints, templates, and one-click ZIP export eliminate redundant engineering"),
    ("🔒 Compliance Risk",  "Automated PII redaction","100%",  ACCENT3,
     "Presidio-powered guardrails eliminate manual data masking effort and audit findings"),
    ("📈 Quality Uplift",   "Pre-deploy simulation",  "3×",   RGBColor(0xA8,0x55,0xF7),
     "Test harness catches hallucinations and regressions before production — 3× fewer incidents"),
    ("🔁 Reuse Rate",       "Shared blueprint library","60%↑", RGBColor(0x0E,0x79,0xB2),
     "Marketplace of proven agent patterns reduces rebuild effort across teams and clients"),
    ("📊 Observability",    "Full trace coverage",    "100%", RGBColor(0xDC,0x26,0x26),
     "OTEL + Audit logs give L1 support full incident context — MTTR reduced by 70%"),
]
for i, (title, metric, value, clr, detail) in enumerate(rois):
    row, col = divmod(i, 3)
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.1 + row * 2.25)
    rect(sl, x, y, Inches(4.0), Inches(2.0), CARD_BG)
    txt(sl, title, x + Inches(0.15), y + Inches(0.12), Inches(3.7), Inches(0.45),
        size=14, bold=True, color=clr)
    txt(sl, value, x + Inches(0.15), y + Inches(0.55), Inches(1.5), Inches(0.7),
        size=36, bold=True, color=WHITE)
    txt(sl, metric, x + Inches(0.15), y + Inches(1.2), Inches(3.7), Inches(0.35),
        size=12, bold=True, color=ACCENT2)
    txt(sl, detail, x + Inches(0.15), y + Inches(1.52), Inches(3.7), Inches(0.42),
        size=10, color=LIGHT_GRAY)
slide_num(sl, 6)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — COMPETITIVE COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Competitive Landscape", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
    size=28, bold=True, color=WHITE)

headers = ["Feature", "AgentForge", "Lyzr Architect", "LangChain", "Microsoft AutoGen", "CrewAI"]
col_w   = [Inches(2.8), Inches(2.0), Inches(1.8), Inches(1.6), Inches(2.0), Inches(1.6)]
col_x   = [Inches(0.35)]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

rows = [
    ["Prompt-to-Agent (NL → Stack)", "✅ Full", "✅ Yes", "❌ No", "❌ No", "❌ No"],
    ["Visual Workflow Builder",       "✅ Yes",  "✅ Yes", "❌ No", "⚡ Basic", "❌ No"],
    ["Built-in RAG Pipeline",         "✅ Yes",  "✅ Yes", "⚡ DIY", "❌ No",  "❌ No"],
    ["PII / Guardrails (Presidio)",   "✅ Yes",  "⚡ Partial","❌ No","❌ No","❌ No"],
    ["Hallucination Detection",       "✅ Yes",  "⚡ Partial","❌ No","❌ No","❌ No"],
    ["Multi-Agent Orchestration",     "✅ Yes",  "✅ Yes", "✅ Yes","✅ Yes","✅ Yes"],
    ["Voice Agent (STT/TTS)",         "✅ Yes",  "❌ No",  "❌ No", "❌ No", "❌ No"],
    ["Control Plane (live modify)",   "✅ Yes",  "⚡ Basic","❌ No","❌ No","❌ No"],
    ["Audit Logs + RBAC",             "✅ Yes",  "⚡ Partial","❌ No","❌ No","❌ No"],
    ["OpenTelemetry Observability",   "✅ Yes",  "❌ No",  "❌ No", "❌ No", "❌ No"],
    ["One-click Deployable ZIP",      "✅ Yes",  "✅ Yes", "❌ No", "❌ No", "❌ No"],
    ["Simulation / Eval Engine",      "✅ Yes",  "⚡ Basic","❌ No","⚡ Basic","❌ No"],
    ["Azure-native (no vendor lock)", "✅ Yes",  "❌ OpenAI","❌ Mix","✅ Azure","❌ Mix"],
    ["On-prem / Private Deployment",  "✅ Yes",  "❌ SaaS", "✅ Yes","✅ Yes","✅ Yes"],
]

row_h = Inches(0.36)
header_y = Inches(1.05)

for ci, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
    fill = ACCENT if ci == 1 else CARD_BG
    rect(sl, x, header_y, w - Inches(0.04), row_h, fill)
    txt(sl, hdr, x + Inches(0.06), header_y + Inches(0.04), w - Inches(0.1), row_h - Inches(0.06),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

for ri, row_data in enumerate(rows):
    y = header_y + row_h + ri * row_h
    fill_row = RGBColor(0x16, 0x20, 0x35) if ri % 2 == 0 else CARD_BG
    for ci, (val, x, w) in enumerate(zip(row_data, col_x, col_w)):
        fill = RGBColor(0x1A, 0x30, 0x50) if ci == 1 else fill_row
        clr  = ACCENT2 if "✅" in val else (ACCENT3 if "⚡" in val else RGBColor(0xEF,0x44,0x44))
        if ci == 0: clr = LIGHT_GRAY
        rect(sl, x, y, w - Inches(0.04), row_h, fill)
        txt(sl, val, x + Inches(0.06), y + Inches(0.04), w - Inches(0.1), row_h - Inches(0.06),
            size=10, color=clr, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
slide_num(sl, 7)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — LEADERSHIP NARRATIVE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Leadership Narrative", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
    size=28, bold=True, color=WHITE)

narrative = [
    ("🎯 Strategic Alignment",
     "AgentForge directly advances Accenture's AI-first mandate by creating a reusable internal platform that "
     "accelerates client delivery of AI agents. Every engagement that uses AgentForge compresses delivery timelines "
     "by 60–85%, improving margin and client NPS simultaneously."),
    ("🏆 Competitive Differentiation",
     "No competitor today combines Prompt-to-Agent generation, built-in enterprise guardrails (PII+hallucination), "
     "voice agents, and a full observability stack in a single deployable platform. AgentForge positions Accenture as "
     "the go-to partner for enterprise-grade AI agent deployments."),
    ("💼 Client Value Proposition",
     "Clients receive a production-ready, governable AI agent platform aligned with their Azure investment. They get "
     "faster innovation cycles, lower risk (compliance-ready guardrails), and full transparency (audit logs, OTEL) — "
     "without needing to build internal AI infrastructure from scratch."),
    ("🌱 Platform Network Effect",
     "The Blueprints Marketplace creates a virtuous cycle: every agent built on AgentForge adds to the shared library, "
     "reducing future build cost. Over 24 months this compounds into a proprietary IP moat for Accenture's AI practice."),
]
for i, (title, body) in enumerate(narrative):
    y = Inches(1.1 + i * 1.45)
    rect(sl, Inches(0.4), y, Inches(0.08), Inches(1.25), ACCENT)
    txt(sl, title, Inches(0.65), y + Inches(0.05), Inches(12), Inches(0.45),
        size=15, bold=True, color=ACCENT)
    txt(sl, body, Inches(0.65), y + Inches(0.48), Inches(12), Inches(0.72),
        size=12.5, color=LIGHT_GRAY)
slide_num(sl, 8)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — KPIs & SUCCESS METRICS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "KPIs & Success Metrics", Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
    size=28, bold=True, color=WHITE)

kpi_groups = [
    ("📐 Platform Adoption", ACCENT, [
        ("Active Users / Month",        "Target: 50+ developers",    "Track: Auth login events"),
        ("Agents Created",              "Target: 200+ agents",       "Track: agents table count"),
        ("Projects / Blueprints Saved", "Target: 100+ projects",     "Track: workflows count"),
    ]),
    ("⚡ Developer Velocity", ACCENT2, [
        ("Time-to-First-Agent",         "Target: <5 minutes",        "Track: Architect session duration"),
        ("ZIP Downloads / Week",        "Target: 20+/week",          "Track: buildSourceZip calls"),
        ("Simulation Pass Rate",        "Target: >90%",              "Track: simulation.status=passed"),
    ]),
    ("🛡️ Governance & Safety", ACCENT3, [
        ("PII Triggers / 1000 calls",   "Target: <5 incidents",      "Track: guardrail_triggered=true"),
        ("Hallucination Rate",          "Target: <2%",               "Track: hallucination_triggered"),
        ("Audit Log Coverage",          "Target: 100%",              "Track: audit_logs completeness"),
    ]),
    ("📊 Reliability & Perf", RGBColor(0xA8,0x55,0xF7), [
        ("API P95 Latency",             "Target: <800ms",            "Track: OTEL latency histograms"),
        ("Workflow Success Rate",       "Target: >98%",              "Track: workflow_runs status"),
        ("System Uptime",               "Target: 99.9%",             "Track: health endpoint + OTEL"),
    ]),
]
for gi, (group_name, clr, items) in enumerate(kpi_groups):
    col = gi % 2
    row = gi // 2
    gx = Inches(0.4 + col * 6.5)
    gy = Inches(1.05 + row * 2.95)
    rect(sl, gx, gy, Inches(6.2), Inches(2.7), CARD_BG)
    rect(sl, gx, gy, Inches(6.2), Inches(0.42), clr)
    txt(sl, group_name, gx + Inches(0.15), gy + Inches(0.06), Inches(6), Inches(0.32),
        size=14, bold=True, color=WHITE)
    for ki, (name, target, track) in enumerate(items):
        ky = gy + Inches(0.52 + ki * 0.7)
        txt(sl, f"• {name}", gx + Inches(0.15), ky, Inches(2.5), Inches(0.35), size=11, bold=True, color=LIGHT_GRAY)
        txt(sl, target, gx + Inches(2.7), ky, Inches(1.8), Inches(0.35), size=11, color=clr)
        txt(sl, track,  gx + Inches(4.5), ky, Inches(1.7), Inches(0.35), size=10, color=RGBColor(0x94,0xA3,0xB8))
slide_num(sl, 9)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — DATABASE SCHEMA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Database Schema (PostgreSQL 16 + pgvector)", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
    size=26, bold=True, color=WHITE)

tables = [
    ("users",          ACCENT,  ["id (PK)", "email (unique)", "hashed_password", "full_name", "role (admin|developer|viewer)", "is_active", "created_at"]),
    ("agents",         ACCENT2, ["id (PK)", "name", "description", "system_prompt", "model (gpt-5.4-mini)", "tools (JSON)", "guardrails (JSON)", "created_by (FK→users)", "agent_type", "current_version"]),
    ("agent_versions", RGBColor(0xA8,0x55,0xF7), ["id (PK)", "agent_id (FK→agents)", "version (int)", "snapshot (JSON)", "created_at"]),
    ("workflows",      ACCENT3, ["id (PK)", "name", "description", "nodes (JSON)", "edges (JSON)", "created_by (FK→users)", "created_at", "updated_at"]),
    ("workflow_runs",  RGBColor(0xDC,0x26,0x26), ["id (PK)", "workflow_id (FK→workflows)", "trigger_input", "final_output", "status", "node_logs (JSON)", "total_duration_ms", "triggered_at"]),
    ("knowledge_bases",RGBColor(0x0E,0x79,0xB2), ["id (PK)", "name", "description", "agent_id (FK→agents)", "created_by", "created_at"]),
    ("documents",      RGBColor(0x06,0xB6,0xD4), ["id (PK)", "kb_id (FK→knowledge_bases)", "filename", "content (TEXT)", "chunk_count", "status (processing|ready)", "created_at"]),
    ("audit_logs",     RGBColor(0xF4,0x3F,0x5E), ["id (PK)", "user_id", "action", "resource_type", "resource_id", "input_snapshot (JSON)", "output_snapshot (JSON)", "guardrail_triggered", "latency_ms", "trace_id", "created_at"]),
    ("api_keys",       RGBColor(0x84,0xCC,0x16), ["id (PK)", "name", "key_prefix", "hashed_key", "user_id (FK→users)", "is_active", "created_at"]),
    ("voice_sessions", RGBColor(0xFD,0xBA,0x74), ["id (PK)", "agent_id (FK→agents)", "user_id", "transcript (JSON)", "status", "created_at"]),
]
col_w = Inches(4.1)
col_gap = Inches(0.1)
for i, (name, clr, fields) in enumerate(tables):
    col = i % 3
    row = i // 3
    x = Inches(0.35 + col * (4.1 + 0.2))
    y = Inches(1.05 + row * 1.65)
    rect(sl, x, y, col_w, Inches(1.5), CARD_BG)
    rect(sl, x, y, col_w, Inches(0.38), clr)
    txt(sl, f"⬛ {name}", x + Inches(0.1), y + Inches(0.04), col_w - Inches(0.15), Inches(0.32),
        size=13, bold=True, color=WHITE)
    body = "  ·  ".join(fields[:4]) + ("\n  ·  " + "  ·  ".join(fields[4:]) if len(fields) > 4 else "")
    txt(sl, body, x + Inches(0.1), y + Inches(0.42), col_w - Inches(0.15), Inches(1.0),
        size=9.5, color=LIGHT_GRAY)
slide_num(sl, 10)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Technology Stack", Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
    size=28, bold=True, color=WHITE)

stacks = [
    ("🖥️ Frontend", ACCENT, [
        "React 18.3 + Vite 5 + TypeScript 5",
        "TailwindCSS 3 (utility-first styling)",
        "React Flow (visual workflow canvas)",
        "TanStack Query v4 (server state)",
        "Lucide React (icon system)",
        "react-hot-toast (notifications)",
        "Axios (API client)",
    ]),
    ("⚙️ Backend", ACCENT2, [
        "FastAPI 0.115 + Uvicorn (ASGI)",
        "Python 3.12",
        "SQLAlchemy 2.0 (async ORM)",
        "Alembic (migrations)",
        "Pydantic v2 (validation)",
        "python-jose (JWT auth)",
        "passlib + bcrypt (password hashing)",
    ]),
    ("🤖 AI / ML", ACCENT3, [
        "Azure OpenAI GPT-5.4-mini + GPT-5.4-mini",
        "LangChain 0.3 + LangChain-OpenAI",
        "OpenAI SDK 1.84",
        "FAISS (vector similarity)",
        "pgvector 0.3 (PostgreSQL vectors)",
        "Presidio (PII detection/anonymization)",
        "Azure AI Search (enterprise RAG)",
    ]),
    ("🗄️ Data & Infra", RGBColor(0xA8,0x55,0xF7), [
        "PostgreSQL 16 + pgvector extension",
        "asyncpg (async DB driver)",
        "Docker + Docker Compose",
        "Azure Speech (STT + TTS)",
        "OpenTelemetry (traces + metrics)",
        "Azure Monitor (OTEL exporter)",
        "pypdf + python-docx (doc parsing)",
    ]),
]
for i, (group, clr, items) in enumerate(stacks):
    x = Inches(0.35 + i * 3.24)
    rect(sl, x, Inches(1.05), Inches(3.1), Inches(5.8), CARD_BG)
    rect(sl, x, Inches(1.05), Inches(3.1), Inches(0.5), clr)
    txt(sl, group, x + Inches(0.12), Inches(1.1), Inches(2.9), Inches(0.42),
        size=15, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(sl, f"▸ {item}", x + Inches(0.12), Inches(1.65 + j * 0.68),
            Inches(2.9), Inches(0.58), size=12, color=LIGHT_GRAY)
slide_num(sl, 11)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — API MODULE MAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "API Module Map  (16 Routers · FastAPI)", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
    size=26, bold=True, color=WHITE)

apis = [
    ("/api/auth",         "JWT login, register, me endpoint"),
    ("/api/agents",       "CRUD agents, versions, tools, guardrails config"),
    ("/api/rag",          "Knowledge bases, document upload & chunking, semantic search"),
    ("/api/tools",        "Tool registry — list, register, invoke native tools"),
    ("/api/builder",      "Visual workflow save/load/run (nodes + edges)"),
    ("/api/simulation",   "Run agent test scenarios, score outputs"),
    ("/api/architect",    "Prompt-to-Agent generation, ZIP export (RAG Scaffold & Custom Code)"),
    ("/api/control-plane","Live agent status, patch, pause, resume running workflows"),
    ("/api/safety",       "Guardrail policy configuration per agent"),
    ("/api/evaluations",  "Batch evaluation runs, scoring metrics, version comparison"),
    ("/api/voice",        "Voice session management, STT→LLM→TTS pipeline"),
    ("/api/api-keys",     "Issue, list, revoke programmatic API keys"),
    ("/api/team",         "User invite, role assignment, team management"),
    ("/api/audit",        "Query audit logs, filter by user/action/resource/date"),
    ("/api/telemetry",    "OpenTelemetry trace ingestion (OTLP/HTTP)"),
    ("/api/health",       "Health check endpoint for load balancers and Kubernetes"),
]
col1 = apis[:8]; col2 = apis[8:]
for i, (route, desc) in enumerate(col1):
    y = Inches(1.05 + i * 0.74)
    rect(sl, Inches(0.35), y, Inches(6.3), Inches(0.65), CARD_BG)
    txt(sl, route, Inches(0.5), y + Inches(0.07), Inches(2.5), Inches(0.5),
        size=12, bold=True, color=ACCENT)
    txt(sl, desc, Inches(3.0), y + Inches(0.07), Inches(3.7), Inches(0.5),
        size=11, color=LIGHT_GRAY)
for i, (route, desc) in enumerate(col2):
    y = Inches(1.05 + i * 0.74)
    rect(sl, Inches(6.95), y, Inches(6.0), Inches(0.65), CARD_BG)
    txt(sl, route, Inches(7.1), y + Inches(0.07), Inches(2.5), Inches(0.5),
        size=12, bold=True, color=ACCENT2)
    txt(sl, desc, Inches(9.6), y + Inches(0.07), Inches(3.3), Inches(0.5),
        size=11, color=LIGHT_GRAY)
slide_num(sl, 12)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — FRONTEND PAGE MAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Frontend Page Map  (24 Pages · React SPA)", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
    size=26, bold=True, color=WHITE)

pages = [
    ("🏠", "Home",               "Landing, quick actions, recent activity"),
    ("🧠", "Planning Architect", "Conversational prompt-to-agent, session history, plan/agents/app/db tabs"),
    ("🔧", "Agent Studio",       "Create, edit, test, deploy, publish agents"),
    ("📐", "Workflow Builder",   "React Flow drag-drop canvas, node config, trigger & run"),
    ("📚", "Knowledge Bases",    "Create KB, upload docs, view chunk status"),
    ("🎮", "Playground",         "Interactive agent testing with streaming responses"),
    ("🧪", "Evaluations",        "Batch test runs, scoring, version comparison"),
    ("🔒", "Safety",             "Guardrail policies per agent, PII/hallucination config"),
    ("📊", "Usage & Traceability","Audit log viewer, latency stats, filter by action/agent"),
    ("🌐", "Control Plane",      "Live workflow status, modify/pause/resume running agents"),
    ("🔭", "Workflow Observability","Per-run execution trace, node-level latency breakdown"),
    ("🗣️", "Voice Agents",       "Voice session management, STT/TTS config, audio test"),
    ("🗂️", "Blueprints",         "Reusable agent blueprints library, clone to project"),
    ("📝", "Prompt Library",     "Curated prompts, search, copy-to-agent"),
    ("💡", "What Should I Build","AI-guided project ideation based on use-case description"),
    ("🗺️", "My Projects",        "Project cards, download ZIP, open in Architect"),
    ("🌍", "Published Projects", "Published agents visible to org members"),
    ("🤝", "Shared Projects",    "Projects shared with you across the org"),
    ("🛒", "Marketplace",        "Community agent templates — import & customise"),
    ("🔑", "API Keys",           "Generate, list, revoke programmatic keys"),
    ("👥", "Team Members",       "Invite, assign roles, deactivate users"),
    ("📅", "Agent Versions",     "Version history, diff, rollback"),
    ("👤", "Profile",            "User profile, password change"),
    ("⚙️", "Settings",           "Org settings, Azure config, notification prefs"),
]
cw = Inches(3.06); ch = Inches(0.62)
for i, (icon, name, desc) in enumerate(pages):
    col = i % 4; row = i // 4
    x = Inches(0.35 + col * 3.24); y = Inches(1.05 + row * 0.68)
    rect(sl, x, y, cw, ch, CARD_BG)
    txt(sl, f"{icon} {name}", x + Inches(0.1), y + Inches(0.04), Inches(1.3), ch - Inches(0.1),
        size=11, bold=True, color=ACCENT)
    txt(sl, desc, x + Inches(1.4), y + Inches(0.07), cw - Inches(1.5), ch - Inches(0.1),
        size=9.5, color=LIGHT_GRAY)
slide_num(sl, 13)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — DEVELOPER SETUP (Part 1: Prerequisites & Backend)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Developer Setup Guide  — Part 1: Prerequisites & Backend", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
    size=22, bold=True, color=WHITE)

prereqs = [
    "Python 3.12+  ·  Node.js 20+ & npm  ·  PostgreSQL 16  ·  Git",
    "pgvector extension: CREATE EXTENSION IF NOT EXISTS vector;",
    "Azure OpenAI resource with gpt-5.4-mini + gpt-5.4-mini deployments",
    "Optional: Azure Speech, Azure AI Search, Azure Monitor",
]
txt(sl, "📋  Prerequisites", Inches(0.5), Inches(0.95), Inches(6), Inches(0.4),
    size=15, bold=True, color=ACCENT)
for i, p in enumerate(prereqs):
    txt(sl, f"  {p}", Inches(0.5), Inches(1.38 + i * 0.32), Inches(12.2), Inches(0.3),
        size=11, color=LIGHT_GRAY)

backend_steps = [
    ("1", "Clone & enter backend",
     "git clone <repo>  &&  cd AgentForge/backend"),
    ("2", "Create virtual environment",
     "python -m venv venv  &&  .\\venv\\Scripts\\activate  (Windows)\nsource venv/bin/activate  (Mac/Linux)"),
    ("3", "Install dependencies",
     "pip install -r requirements.txt"),
    ("4", "Configure environment",
     "Copy .env.example → .env\nSet: AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_API_KEY,\n     AZURE_OPENAI_DEPLOYMENT_GPT4O=gpt-5.4-mini,\n     DATABASE_URL=postgresql+asyncpg://postgres:postgres@localhost:5432/agentforge"),
    ("5", "Create database",
     "createdb agentforge  (PostgreSQL CLI)\npsql -d agentforge -c 'CREATE EXTENSION IF NOT EXISTS vector;'"),
    ("6", "Run the API server",
     "uvicorn app.main:app --reload --port 8000\n✅  API: http://localhost:8000\n✅  Swagger: http://localhost:8000/docs"),
]
txt(sl, "⚙️  Backend Setup", Inches(0.5), Inches(2.72), Inches(6), Inches(0.4),
    size=15, bold=True, color=ACCENT2)
for step in backend_steps:
    num, title, cmd = step
    idx = backend_steps.index(step)
    y = Inches(3.15 + idx * 0.67)
    chip(sl, num, Inches(0.5), y + Inches(0.04), Inches(0.32), Inches(0.32), ACCENT2, 11)
    txt(sl, title, Inches(0.92), y + Inches(0.04), Inches(2.5), Inches(0.32), size=12, bold=True, color=WHITE)
    txt(sl, cmd, Inches(3.5), y, Inches(9.2), Inches(0.6), size=10, color=RGBColor(0x86,0xEF,0xAC))
slide_num(sl, 14)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — DEVELOPER SETUP (Part 2: Frontend & Ports)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Developer Setup Guide  — Part 2: Frontend, Ports & Verification", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
    size=22, bold=True, color=WHITE)

frontend_steps = [
    ("1", "Enter frontend dir",  "cd AgentForge/frontend"),
    ("2", "Install packages",    "npm install"),
    ("3", "Start dev server",    "npm run dev\n✅  Frontend: http://localhost:5173"),
    ("4", "Login credentials",   "Admin: admin@agentforge.io / admin123\n(seeded automatically on first startup)"),
]
txt(sl, "🖥️  Frontend Setup", Inches(0.5), Inches(1.0), Inches(6), Inches(0.4),
    size=15, bold=True, color=ACCENT)
for i, (num, title, cmd) in enumerate(frontend_steps):
    y = Inches(1.45 + i * 0.72)
    chip(sl, num, Inches(0.5), y + Inches(0.04), Inches(0.32), Inches(0.32), ACCENT, 11)
    txt(sl, title, Inches(0.92), y + Inches(0.04), Inches(2.5), Inches(0.32), size=12, bold=True, color=WHITE)
    txt(sl, cmd, Inches(3.5), y, Inches(9.2), Inches(0.62), size=10, color=RGBColor(0x86,0xEF,0xAC))

# Port table
txt(sl, "🔌  Port Assignment", Inches(0.5), Inches(4.0), Inches(6), Inches(0.4),
    size=15, bold=True, color=ACCENT3)
ports = [
    ("8000", "AgentForge Backend (FastAPI)",              "Main API, Swagger at /docs"),
    ("5173", "AgentForge Frontend (Vite)",                "React SPA dev server"),
    ("8001", "RAG Scaffold Backend (downloaded ZIP)",     "Routes: /api/chat, /api/documents"),
    ("8002", "Custom Code Backend (downloaded ZIP)",      "Routes: /api/ask-question, /api/upload-document"),
    ("5176", "RAG Scaffold Frontend (downloaded ZIP)",    "npm run dev in extracted folder"),
    ("5177", "Custom Code Frontend (downloaded ZIP)",     "npm run dev in extracted folder"),
    ("4318", "OpenTelemetry OTLP/HTTP (optional)",        "Jaeger / Grafana OTEL collector"),
]
ph = Inches(0.42)
for i, (port, service, note) in enumerate(ports):
    y = Inches(4.5 + i * ph)
    fill = CARD_BG if i % 2 == 0 else RGBColor(0x16,0x20,0x35)
    rect(sl, Inches(0.35), y, Inches(12.6), ph, fill)
    txt(sl, port, Inches(0.5), y + Inches(0.05), Inches(0.7), ph - Inches(0.08),
        size=12, bold=True, color=ACCENT3)
    txt(sl, service, Inches(1.3), y + Inches(0.05), Inches(5.5), ph - Inches(0.08),
        size=11, color=WHITE)
    txt(sl, note, Inches(7.0), y + Inches(0.05), Inches(5.8), ph - Inches(0.08),
        size=10, color=LIGHT_GRAY)
slide_num(sl, 15)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — PLANNING ARCHITECT DEEP DIVE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Planning Architect — Prompt-to-Agent Deep Dive", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
    size=24, bold=True, color=WHITE)

steps = [
    ("1", "User describes use case in plain English",
     "e.g. 'Build a policy analysis agent that ingests company documents and answers compliance questions'"),
    ("2", "GPT-5.4-mini generates full project plan",
     "Architecture · Tech stack · Agent configs · Tool list · DB schema · Security considerations"),
    ("3", "Agents tab — agent cards generated",
     "Each agent gets: name, type, model, system prompt, tools list, guardrails config"),
    ("4", "App tab — live sandbox preview",
     "Interactive 3-panel chat UI rendered immediately: left sidebar (docs), main chat, right panel (KB stats)"),
    ("5", "Download RAG Scaffold ZIP",
     "Complete React + FastAPI project using /api/chat + /api/documents routes → proxy target port 8001"),
    ("6", "Download Custom Code ZIP",
     "GPT-5.4-mini generates bespoke App.tsx → renderMarkdown + /api/ask-question + /api/upload-document → port 8002"),
    ("7", "Database tab — entity diagram",
     "Auto-generated PostgreSQL schema matching the planned agent's data requirements"),
]
for i, (num, title, detail) in enumerate(steps):
    row = i % 4; col = i // 4
    y = Inches(1.05 + row * 1.55)
    x = Inches(0.35 + col * 6.55)
    rect(sl, x, y, Inches(6.3), Inches(1.4), CARD_BG)
    chip(sl, num, x + Inches(0.15), y + Inches(0.15), Inches(0.38), Inches(0.38), ACCENT, 13)
    txt(sl, title, x + Inches(0.65), y + Inches(0.12), Inches(5.5), Inches(0.45),
        size=13, bold=True, color=WHITE)
    txt(sl, detail, x + Inches(0.15), y + Inches(0.65), Inches(6.0), Inches(0.65),
        size=11, color=LIGHT_GRAY)
slide_num(sl, 16)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — GUARDRAILS & SAFETY ENGINE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Enterprise Guardrails & Safety Engine", Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
    size=26, bold=True, color=WHITE)

# Flow
flow_items = [
    ("User Input", Inches(0.4)),
    ("Guardrails Check", Inches(2.9)),
    ("PII Redaction\n(Presidio)", Inches(5.4)),
    ("Hallucination\nDetection", Inches(7.9)),
    ("Agent Response\n(Sanitized)", Inches(10.4)),
]
for i, (label, x) in enumerate(flow_items):
    clr = ACCENT if i == 0 else (ACCENT2 if i == 4 else ACCENT3)
    rect(sl, x, Inches(1.2), Inches(2.2), Inches(0.85), clr)
    txt(sl, label, x + Inches(0.05), Inches(1.25), Inches(2.1), Inches(0.75),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, "→", x + Inches(2.2), Inches(1.45), Inches(0.5), Inches(0.4),
            size=22, bold=True, color=ACCENT)

details = [
    ("🔍 PII Entities Detected", ACCENT3, [
        "EMAIL_ADDRESS — e.g. john@example.com → <EMAIL>",
        "PHONE_NUMBER — e.g. 555-123-4567 → <PHONE_NUMBER>",
        "PERSON — named entity recognition",
        "CREDIT_CARD — 16-digit card patterns",
        "US_SSN — Social Security Numbers",
        "Powered by Microsoft Presidio Analyzer + Anonymizer",
    ]),
    ("⚠️ Hallucination Phrases", ACCENT, [
        "Uncertainty phrase matching on agent output",
        "Triggers: 'I'm not sure but', 'I think maybe'",
        "Triggers: 'I cannot verify', 'it might be'",
        "Triggers: 'I believe but I'm not certain'",
        "Result: flagged in audit log + response annotated",
        "Per-agent enable/disable via guardrails config",
    ]),
    ("📋 Audit Trail", ACCENT2, [
        "Every guardrail trigger logged to audit_logs table",
        "Fields: guardrail_triggered (bool), trace_id, latency_ms",
        "Searchable via /api/audit with filter params",
        "Exportable for compliance reporting",
        "OpenTelemetry span attributes: pii_triggered, hallucination_triggered",
        "Zero-trust: logs never omitted even on error",
    ]),
]
for i, (title, clr, items) in enumerate(details):
    x = Inches(0.35 + i * 4.35)
    y = Inches(2.4)
    rect(sl, x, y, Inches(4.1), Inches(4.5), CARD_BG)
    rect(sl, x, y, Inches(4.1), Inches(0.45), clr)
    txt(sl, title, x + Inches(0.12), y + Inches(0.06), Inches(4.0), Inches(0.35),
        size=13, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(sl, f"• {item}", x + Inches(0.12), y + Inches(0.56 + j * 0.63),
            Inches(3.9), Inches(0.55), size=11, color=LIGHT_GRAY)
slide_num(sl, 17)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Product Roadmap", Inches(0.5), Inches(0.3), Inches(8), Inches(0.6),
    size=28, bold=True, color=WHITE)

phases = [
    ("✅ Phase 1 — Complete  (Q1 2026)", ACCENT2, [
        "Core platform: Auth + RBAC + Agent CRUD",
        "Azure OpenAI integration (GPT-5.4-mini + GPT-5.4-mini)",
        "RAG pipeline (FAISS + pgvector)",
        "Guardrails engine (PII + hallucination)",
        "Multi-agent orchestrator",
        "Planning Architect (Prompt-to-Agent + ZIP export)",
        "OpenTelemetry observability + Audit logs",
    ]),
    ("🚀 Phase 2 — In Progress  (Q2–Q3 2026)", ACCENT3, [
        "Visual Workflow Builder (React Flow canvas)",
        "Agent Simulation & Evaluation engine",
        "Voice Agent pipeline (Azure Speech STT/TTS)",
        "Safety policy configuration UI",
        "Blueprints Marketplace (shared templates)",
        "Team management + API Key management",
        "Control Plane (live agent modification)",
    ]),
    ("🔮 Phase 3 — Planned  (Q4 2026)", RGBColor(0xA8,0x55,0xF7), [
        "Multi-modal agents (image + document input)",
        "Agent-to-Agent collaboration (A2A protocol)",
        "Kubernetes Helm chart deployment",
        "SSO / Azure AD integration",
        "Fine-tuning pipeline for domain adaptation",
        "Cost analytics per agent (token usage billing)",
        "Agent performance benchmarking dashboard",
    ]),
    ("💡 Phase 4 — Vision  (2027+)", RGBColor(0xFD,0xBA,0x74), [
        "Autonomous agent self-improvement loop",
        "Cross-org agent marketplace with licensing",
        "Real-time multi-agent collaboration canvas",
        "Regulatory compliance packs (HIPAA, GDPR, SOC2)",
        "Edge deployment (on-device agents)",
        "Federated learning for privacy-preserving training",
        "GenAI Operations (GenAIOps) full lifecycle",
    ]),
]
for i, (phase, clr, items) in enumerate(phases):
    x = Inches(0.35 + i * 3.24)
    rect(sl, x, Inches(1.05), Inches(3.1), Inches(5.8), CARD_BG)
    rect(sl, x, Inches(1.05), Inches(3.1), Inches(0.52), clr)
    txt(sl, phase, x + Inches(0.1), Inches(1.1), Inches(3.0), Inches(0.42),
        size=11, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(sl, f"▸ {item}", x + Inches(0.1), Inches(1.68 + j * 0.68),
            Inches(3.0), Inches(0.6), size=11, color=LIGHT_GRAY)
slide_num(sl, 18)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — ENTERPRISE HARDENING PASS — OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Enterprise Hardening Pass — What Changed", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
    size=26, bold=True, color=WHITE)
txt(sl, "The Architect (Prompt-to-Agent) flow was taken from a working prototype to a "
    "production-grade generator through five focused workstreams — each independently "
    "designed, implemented, live-tested against real Azure OpenAI calls, and code-reviewed.",
    Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.6), size=13, color=LIGHT_GRAY)

hardening = [
    ("🎯", "Domain-Aware Generation",
     "Expanded from 2 hardcoded app types to 10 real domains (HR, Sales, Legal, Support, Marketing, "
     "Dev Tool, Analyst, Data, Chatbot, Custom) — each with its own enterprise-grade layout, labels, and 46 upgraded multi-agent prompts."),
    ("📊", "Real Data, Not Fabricated Demos",
     "Sandbox previews and downloadable Custom Code now seed dashboards/tables from real uploaded sample "
     "data (CSV/Excel) instead of GPT-4o inventing plausible-looking numbers — verified exact-match against source files."),
    ("🗄️", "Enterprise Layers 3–5",
     "Auto-generated PostgreSQL schema + migrations on first run · OpenTelemetry tracing on every LLM call · "
     "an AI-scored feedback loop that few-shots top-rated plans back into future generations."),
    ("🔐", "Real SSO Integration",
     "Azure AD / Entra ID keyword-detected from the plan → genuine JWT-validation middleware (JWKS, aud/iss/exp "
     "checks) + real MSAL frontend login in the Custom Code ZIP, gated by an SSO_ENABLED toggle for local dev."),
    ("📁", "Real File Upload Parsing",
     "Sandbox: real client-side CSV/XLSX parsing (no fake success toasts). Custom Code backend: real PDF "
     "(PyPDF2) and DOCX (python-docx) extraction, closing a gap where both libraries were required but never used."),
]
for i, (icon, title, body) in enumerate(hardening):
    y = Inches(1.68 + i * 0.98)
    rect(sl, Inches(0.4), y, Inches(0.08), Inches(0.85), ACCENT)
    txt(sl, f"{icon}  {title}", Inches(0.65), y + Inches(0.02), Inches(12), Inches(0.4),
        size=14, bold=True, color=ACCENT)
    txt(sl, body, Inches(0.65), y + Inches(0.42), Inches(12.1), Inches(0.5),
        size=11, color=LIGHT_GRAY)
slide_num(sl, 19)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — DOMAIN-AWARE GENERATION & REAL DATA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Domain-Aware Generation Across 10 Categories", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
    size=24, bold=True, color=WHITE)

domains = ["General", "Marketing", "Sales", "Legal", "HR", "Support",
           "Productivity", "Development", "Analysts", "Data & Analysis"]
dcol_w = Inches(2.4); dcol_h = Inches(0.55)
for i, d in enumerate(domains):
    col = i % 5; row = i // 5
    x = Inches(0.4 + col * 2.5); y = Inches(1.0 + row * 0.65)
    chip(sl, d, x, y, dcol_w, dcol_h, ACCENT if row == 0 else ACCENT2, 12)

txt(sl, "Each domain gets its own enterprise-standard layout (correct sidebar labels, KPI tiles, "
    "charts) instead of a generic one-size-fits-all template — 46 Prompt Library prompts upgraded "
    "to multi-agent specs with a matching realistic sample data file per prompt.",
    Inches(0.4), Inches(2.45), Inches(12.4), Inches(0.65), size=12, color=LIGHT_GRAY)

txt(sl, "Before → After: Real Data Seeding", Inches(0.4), Inches(3.3), Inches(8), Inches(0.4),
    size=15, bold=True, color=ACCENT3)

before_after = [
    ("❌ Before", RGBColor(0xDC,0x26,0x26), [
        "\"Contracts in Review: 128\" — a fabricated round number,",
        "unrelated to the real uploaded contract data",
        "\"Outlook Reminders Sent: 12\" — invented for every render",
        "Same issue in the Custom Code ZIP's seed data",
    ]),
    ("✅ After", ACCENT2, [
        "\"Contracts in Review: 4\" · \"High Risk: 3\" — verified exact",
        "match against the real 10-row sample CSV's actual counts",
        "\"Total contract value: $9,579,000\" — sum of real value_usd",
        "column, byte-for-byte traceable to the uploaded file",
    ]),
]
for i, (title, clr, items) in enumerate(before_after):
    x = Inches(0.4 + i * 6.35)
    rect(sl, x, Inches(3.75), Inches(6.1), Inches(2.9), CARD_BG)
    rect(sl, x, Inches(3.75), Inches(6.1), Inches(0.45), clr)
    txt(sl, title, x + Inches(0.15), Inches(3.8), Inches(5.8), Inches(0.35),
        size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(sl, item, x + Inches(0.15), Inches(4.35 + j * 0.55), Inches(5.8), Inches(0.5),
            size=11, color=LIGHT_GRAY)
slide_num(sl, 20)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — ENTERPRISE LAYERS 3–5
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Enterprise Layers 3–5: DB · Observability · AI Reinforcement", Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
    size=23, bold=True, color=WHITE)

layer_cards = [
    ("🗄️", "Layer 3 — DB Auto-Setup", ACCENT, [
        "Domain-specific PostgreSQL schema (db/init.sql)",
        "auto-generated for each of the 10 app domains",
        "Idempotent run_migrations.py runs on first backend",
        "startup — no manual DB scripting required",
    ]),
    ("📡", "Layer 4 — OpenTelemetry", ACCENT2, [
        "Every LLM call (KB extraction, UI generation,",
        "frontend/backend generation) wrapped in a real span",
        "telemetry.py + docker-compose.jaeger.yml auto-injected",
        "into every downloaded Custom Code ZIP",
    ]),
    ("🔁", "Layer 5 — AI Reinforcement Loop", ACCENT3, [
        "👍 / 👎 feedback widget on every generated plan",
        "GPT-4o self-scores plans across 5 quality dimensions",
        "Top-rated plans few-shot-injected into future generation",
        "prompts — the system improves from its own best output",
    ]),
]
for i, (icon, title, clr, items) in enumerate(layer_cards):
    x = Inches(0.35 + i * 4.24)
    rect(sl, x, Inches(1.1), Inches(4.05), Inches(4.6), CARD_BG)
    rect(sl, x, Inches(1.1), Inches(4.05), Inches(0.55), clr)
    txt(sl, f"{icon}  {title}", x + Inches(0.15), Inches(1.2), Inches(3.8), Inches(0.4),
        size=13, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(sl, item, x + Inches(0.15), Inches(1.85 + j * 0.62), Inches(3.8), Inches(0.55),
            size=11, color=LIGHT_GRAY)

txt(sl, "All three layers were verified live: real docker-compose migrations, real span data visible "
    "in a tracing backend, and a real thumbs-up → /feedback → /feedback/top → few-shot round-trip.",
    Inches(0.4), Inches(5.95), Inches(12.4), Inches(0.5), size=11, color=LIGHT_GRAY)
slide_num(sl, 21)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 22 — REAL SSO + REAL FILE UPLOAD
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Real SSO Integration & Real File Upload Parsing", Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
    size=23, bold=True, color=WHITE)

sso_col = [
    "Plan answers like \"Internal enterprise app with SSO\" used to",
    "only affect descriptive text — zero real auth code was generated",
    "",
    "Now: keyword-detected from the plan summary (same pattern as",
    "domain detection) → triggers a real backend + frontend scaffold:",
    "",
    "• backend/app/auth/sso.py — real JWT validation via Azure AD's",
    "  JWKS endpoint, checks aud/iss/exp, python-jose",
    "• src/auth/msalConfig.ts + useAuth.ts — real @azure/msal-browser",
    "  login flow, attaches real Bearer tokens to API calls",
    "• SSO_ENABLED=false by default — app still runs locally without",
    "  a real Azure AD tenant; flip to true once configured",
]
upload_col = [
    "Upload buttons in both the sandbox and downloaded backend were",
    "decorative — a canned success toast regardless of what was picked",
    "",
    "Now — Sandbox (client-side, no backend):",
    "• Real <input type=\"file\"> + FileReader + the already-loaded",
    "  SheetJS library — genuine CSV/XLSX parsing, real row preview",
    "",
    "Now — Custom Code backend (real Python):",
    "• PyPDF2 for real PDF text extraction, python-docx for real DOCX",
    "• Forbids the \"decode raw bytes as UTF-8\" anti-pattern found",
    "  live in a custom-named ingest endpoint — would have corrupted",
    "  any real PDF/DOCX upload into garbage text",
]
cols = [("🔐 Real SSO Integration", ACCENT, sso_col), ("📁 Real File Upload Parsing", ACCENT2, upload_col)]
for i, (title, clr, lines) in enumerate(cols):
    x = Inches(0.35 + i * 6.35)
    rect(sl, x, Inches(1.05), Inches(6.1), Inches(5.6), CARD_BG)
    rect(sl, x, Inches(1.05), Inches(6.1), Inches(0.5), clr)
    txt(sl, title, x + Inches(0.15), Inches(1.12), Inches(5.8), Inches(0.4),
        size=14, bold=True, color=WHITE)
    body = "\n".join(lines)
    txt(sl, body, x + Inches(0.15), Inches(1.65), Inches(5.85), Inches(4.9), size=10.5, color=LIGHT_GRAY)
slide_num(sl, 22)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 23 — HARDENING ROI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Hardening ROI — Risk & Rework Avoided", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
    size=27, bold=True, color=WHITE)

hroi = [
    ("🎭 Demo Credibility", "Fabricated → real dashboard data", "100%", ACCENT,
     "Client-facing demos now show numbers traceable to real uploaded data, eliminating a\ncredibility risk when a stakeholder cross-checks a KPI against source data"),
    ("🔓 Security Exposure", "SSO answer → real auth code", "0→1", ACCENT3,
     "Closes a real gap where an \"Enterprise app with SSO\" answer previously shipped zero\nauthentication code — a downloaded project would have been deployed wide open"),
    ("📉 Data-Corruption Risk", "PDF/DOCX handling", "100%", RGBColor(0xDC,0x26,0x26),
     "Forbids a found anti-pattern (raw UTF-8 decode on binary files) that would silently\ncorrupt any real PDF/DOCX contract, resume, or document a user actually uploads"),
    ("🧪 Escaped Defects", "Found via live testing, not assumed", "9", RGBColor(0xA8,0x55,0xF7),
     "Nine real bugs (session pollution, stale feedback state, false-positive SSO detection,\nenv-file clobbering, and more) caught by regenerating real apps end-to-end, not just code review"),
    ("🌐 Domain Coverage", "2 hardcoded types → 10 real domains", "5×", ACCENT2,
     "Every Prompt Library category now gets a correct, enterprise-appropriate layout instead\nof a generic template mismatched to the actual business domain"),
    ("🔁 Continuous Improvement", "Manual tuning → self-scoring loop", "∞", RGBColor(0x0E,0x79,0xB2),
     "The AI reinforcement loop means generation quality compounds over time from real\nusage feedback, instead of staying frozen at whatever quality shipped on day one"),
]
for i, (title, metric, value, clr, detail) in enumerate(hroi):
    row, col = divmod(i, 3)
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.05 + row * 2.95)
    rect(sl, x, y, Inches(4.05), Inches(2.75), CARD_BG)
    txt(sl, title, x + Inches(0.15), y + Inches(0.12), Inches(3.75), Inches(0.45),
        size=13, bold=True, color=clr)
    txt(sl, value, x + Inches(0.15), y + Inches(0.55), Inches(1.6), Inches(0.7),
        size=30, bold=True, color=WHITE)
    txt(sl, metric, x + Inches(0.15), y + Inches(1.18), Inches(3.75), Inches(0.35),
        size=11, bold=True, color=ACCENT2)
    txt(sl, detail, x + Inches(0.15), y + Inches(1.52), Inches(3.8), Inches(1.15),
        size=9.5, color=LIGHT_GRAY)
slide_num(sl, 23)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 24 — WORKFLOW BUILDER: CONDITIONAL LOGIC, ROUTING & INTEGRATIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Visual Workflow Builder — Conditional Logic, Routing & Integrations", Inches(0.5), Inches(0.3), Inches(12.4), Inches(0.6),
    size=22, bold=True, color=WHITE)
txt(sl, "A real decision-making execution engine — not just a flowchart. Every branch below "
    "genuinely changes which nodes run, verified against live Azure OpenAI runs.",
    Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.5), size=12.5, color=LIGHT_GRAY)

wf_features = [
    ("❓", "Condition Nodes",
     "Rule (e.g. \"days <= 2\") evaluated safely via simpleeval — never Python's eval() — against variables an LLM extracts from the running text. Only the matching true/false branch executes; the other is fully skipped, including convergent downstream nodes."),
    ("🔀", "Router Nodes — Real Branching",
     "A router's LLM decision is classified against its labeled outgoing edges (e.g. \"Fast\"/\"Deep\") and only ONE branch runs — replacing the previous behavior where every downstream node executed regardless of the routing decision."),
    ("✉️", "Human-in-the-Loop Approval",
     "An approval node genuinely pauses the pipeline, sends a real email via SMTP, and exposes a dedicated /approvals/{run_id} page where a reviewer can Approve or Reject — resuming execution from exactly where it paused."),
    ("🌐", "HTTP Request Node — Outbound API Calls",
     "New node type: configurable method (GET/POST/PUT/PATCH/DELETE), URL, JSON headers and body, with a {{input}} placeholder to inject the previous node's output — giving workflows real external-API integration for the first time."),
    ("🐍", "Faithful Python Export",
     "\"Export Code\" now generates a script that ports the exact same topological sort, condition/approval/router/http_request logic as the live engine — not a flat linear stub — plus an optional --openai flag to run it against a real LLM."),
    ("📄", "Export / Import JSON",
     "A lossless, round-trippable workflow backup format — separate from the one-way Python export — with a genuine \"browse and load\" file picker to restore a canvas exactly as it was, verified after a full wipe-and-reload cycle."),
    ("📊", "Workflow Runs — Approval Tracking",
     "The observability dashboard now surfaces an \"Awaiting Approval\" count and a direct \"Review →\" link on paused runs — closing a gap where there was previously no in-app path to the approval screen at all."),
    ("🎨", "Execution-State Visualization",
     "Nodes now genuinely pulse amber while running and turn green on completion — the color-coding logic existed in code but was never actually wired to canvas-created nodes until this pass; idle edges are blue, executed edges turn green."),
]
for i, (icon, title, body) in enumerate(wf_features):
    row, col = divmod(i, 2)
    x = Inches(0.35 + col * 6.35)
    y = Inches(1.5 + row * 1.42)
    card(sl, x, y, Inches(6.15), Inches(1.32), title, [body], ACCENT, icon)
slide_num(sl, 24)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 25 — KNOWN LIMITATIONS & NEXT HARDENING TARGETS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl); accent_bar(sl)
txt(sl, "Known Limitations & Next Hardening Targets", Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
    size=25, bold=True, color=WHITE)
txt(sl, "Shipping fast means shipping honestly. These are confirmed, currently-open gaps — not "
    "hypothetical risks — found and verified during this build cycle.",
    Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), size=12.5, color=LIGHT_GRAY)

gaps = [
    ("🔧", "Agent Studio Tools Are Not Wired to Execution",
     "The tool checkboxes on an Agent (Slack, GitHub, Jira, Google Drive, web search, calculator) are saved to the database but never read at runtime — checking a box today has zero effect. No LLM function/tool-calling is implemented anywhere in the platform yet."),
    ("🔓", "HTTP Request Node Has No SSRF Guardrails",
     "The new outbound-call node will hit any URL it's given, including internal/private network addresses. Needs a host allowlist or denylist (and ideally a per-org config) before this is safe to expose in a production, multi-tenant deployment."),
    ("⏰", "JWT Sessions Expire Without Silent Refresh",
     "The 8-hour token lifetime has no refresh flow — long sessions eventually fail with a raw \"Invalid token\" error on any authenticated action until the user logs in again. A refresh-token flow would remove this rough edge."),
    ("🗄️", "No Formal Database Migration Tooling",
     "New WorkflowRun columns added this cycle required a manual ALTER TABLE fix after schema drift was hit live. Alembic (already a listed dependency) isn't actually wired into the startup path — schema changes are still ad-hoc."),
    ("🧭", "Exported Python Uses Canvas Order as a Tiebreaker",
     "The topological sort ported into Export Code is correct for the graphs tested, but canvas creation order (not a formal secondary sort key) breaks ties — fine in practice, but worth a closer look for pathological or cyclic graphs."),
]
for i, (icon, title, body) in enumerate(gaps):
    y = Inches(1.68 + i * 0.98)
    rect(sl, Inches(0.4), y, Inches(0.08), Inches(0.88), ACCENT3)
    txt(sl, f"{icon}  {title}", Inches(0.65), y + Inches(0.02), Inches(12), Inches(0.4),
        size=13.5, bold=True, color=ACCENT3)
    txt(sl, body, Inches(0.65), y + Inches(0.42), Inches(12.1), Inches(0.52),
        size=10.5, color=LIGHT_GRAY)
slide_num(sl, 25)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 26 — CLOSING / CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
accent_bar(sl, Inches(0), ACCENT)

rect(sl, 0, Inches(1.5), W, Inches(4.5), RGBColor(0x1A, 0x24, 0x3E))

txt(sl, "AgentForge", Inches(0.5), Inches(1.8), Inches(12), Inches(1.0),
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Build the next generation of enterprise AI agents — today.",
    Inches(0.5), Inches(2.85), Inches(12), Inches(0.8),
    size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

ctas = [
    ("🚀  Start Building", "localhost:5173/architect", ACCENT),
    ("📖  API Docs",        "localhost:8000/docs",      ACCENT2),
    ("📦  Download ZIP",   "Architect → App → Download", ACCENT3),
]
for i, (label, url, clr) in enumerate(ctas):
    x = Inches(1.5 + i * 3.5)
    rect(sl, x, Inches(4.0), Inches(3.1), Inches(0.85), clr)
    txt(sl, label, x, Inches(4.05), Inches(3.1), Inches(0.42),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, url, x, Inches(4.47), Inches(3.1), Inches(0.35),
        size=11, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Accenture Technology  |  n.sureshmanikandan@accenture.com  |  Confidential",
    Inches(0.5), Inches(6.3), Inches(12), Inches(0.4),
    size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
slide_num(sl, 26)


# ── Save ──────────────────────────────────────────────────────────────────────
out = r"C:\Users\n.sureshmanikandan\Repo1\AgentForge\docs\AgentForge-PresentationV7.0.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
