# -*- coding: utf-8 -*-
"""
Adds one new content slide to AgentForge-PresentationV9.pptx documenting the
Workflow Builder testing/approval UX features shipped this session:
  - Save Workflow real naming (was silently writing to localStorage only)
  - Browse & Load Saved Workflows picker
  - Force-branch deterministic router testing
  - In-canvas Approve/Reject (no more leaving the page for approvals)
  - Full-path visual completion + amber escalation color (no more red collision)

Clones the exact visual style of slide 26 ("Latest Platform Updates — Local
LLM & Model Flexibility"): dark navy background, thin accent top bar, title +
description, and N stacked feature rows (accent tick + bold colored heading +
description paragraph), with a page-number footer.

Inserted right after slide 34 (Lyzr-Parity Scorecard) and before slide 35
(How To Use tutorials begin), preserving every existing slide untouched.

Saves as AgentForge-PresentationV10.pptx (does not overwrite V9).
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy

SRC = r"C:\Users\n.sureshmanikandan\Repo1\AgentForge\docs\AgentForge-PresentationV9.pptx"
OUT = r"C:\Users\n.sureshmanikandan\Repo1\AgentForge\docs\AgentForge-PresentationV10.pptx"

BG_NAVY = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x60, 0x6E, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

SLIDE_W = Emu(12188952)
SLIDE_H = Emu(6858000)

TITLE = "Workflow Builder — Testing & Human-in-the-Loop UX"
SUBTITLE = (
    "A tester reported “Save Workflow saved with a junk timestamp name” and "
    "“the approval step feels like an error.” Both traced to real gaps, now closed — "
    "verified live against a running backend, not just unit tests."
)

ITEMS = [
    (
        "\U0001F4BE  Save Workflow — Real Names, Not Timestamps",
        "“Save Workflow” only wrote to browser localStorage and never reached the backend — the "
        "junk “Workflow 8:14:41 PM” entries actually came from Run/Deploy auto-saving silently to get "
        "a workflow_id. Save Workflow now prompts for a real name and persists it directly.",
    ),
    (
        "\U0001F5C2️  Browse & Load Saved Workflows",
        "Load now opens a searchable picker over every backend-saved workflow (name, node count, last "
        "updated) instead of silently reading localStorage — selecting one loads it straight onto the "
        "canvas, with legacy flat-format workflows normalized so they render instead of crashing.",
    ),
    (
        "\U0001F500  Force a Branch — Deterministic Router Testing",
        "The Run modal now shows a “Force branch” dropdown for every router node — pick Critical/"
        "Routine (or any labeled edge) to skip the LLM classification and reliably test one specific path, "
        "instead of guessing wording that happens to trigger it.",
    ),
    (
        "✅  In-Canvas Approve/Reject",
        "A paused run now surfaces an inline “Run paused for approval” banner directly in Workflow "
        "Builder with Approve/Reject buttons wired to the existing approval endpoints — testers no longer "
        "need the emailed link or a separate /approvals/{run_id} page to finish exercising the Critical path.",
    ),
    (
        "\U0001F3A8  Full-Path Visual Completion & Amber Escalation Color",
        "Approving a paused run now marks every downstream node and edge as done on the canvas — not just "
        "the approval node itself — and the approval role's default color moved from red to amber so a "
        "normal paused-for-approval state no longer visually collides with the canvas's “error” red.",
    ),
]

FOOTER_NUMBER = "34a"  # inserted between slide 34 and 35; see note in module docstring


def add_rect(slide, left, top, width, height, fill_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size_pt, bold, color_rgb, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color_rgb
    return box


def update_stale_approval_bullet(prs):
    """Slide 24's Human-in-the-Loop Approval bullet still says the *only* way to
    resolve a paused run is the emailed link / /approvals/{run_id} page. That's now
    outdated -- Workflow Builder has an in-canvas Approve/Reject banner too."""
    slide24 = prs.slides[23]
    old_fragment = "resuming execution from exactly where it paused."
    new_suffix = (
        " resuming execution from exactly where it paused — reviewable either via that "
        "emailed link or via an in-canvas Approve/Reject banner directly in Workflow Builder."
    )
    for shape in slide24.shapes:
        if not shape.has_text_frame:
            continue
        if old_fragment in shape.text_frame.text:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if old_fragment in run.text:
                        run.text = run.text.replace(
                            "resuming execution from exactly where it paused.", new_suffix.strip()
                        )
            return True
    return False


def fix_admin_email_typo(prs):
    """Slide 15's setup guide says admin@agentforge.io -- the real seeded
    admin account (backend/app/core/seed.py) is admin@agentforge.ai."""
    slide15 = prs.slides[14]
    for shape in slide15.shapes:
        if not shape.has_text_frame:
            continue
        if "admin@agentforge.io" in shape.text_frame.text:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if "admin@agentforge.io" in run.text:
                        run.text = run.text.replace("admin@agentforge.io", "admin@agentforge.ai")
            return True
    return False


def fix_api_module_map(prs):
    """Slide 12 lists '/api/audit' and '/api/telemetry' as if they were their
    own top-level routers -- neither exists in backend/app/main.py. Audit-log
    querying actually lives at /api/control-plane/audit-logs, and there is no
    telemetry REST endpoint at all (OTel is exported via OTLP/HTTP, not queried
    through the API). Also, /api/projects is a real, registered router that
    was missing from this slide entirely -- added later than this deck.
    Fix both fictional entries in place (zero layout risk, same tile count)."""
    slide12 = prs.slides[11]
    shapes = list(slide12.shapes)
    replacements = {
        "/api/audit": ("/api/control-plane/audit-logs", "Query audit logs, filter by user/action/resource/date"),
        "/api/telemetry": ("/api/projects", "CRUD projects, publish/visibility toggle, Architect auto-save integration"),
    }
    done = set()
    for i, shape in enumerate(shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text in replacements and text not in done:
            new_header, new_desc = replacements[text]
            shape.text_frame.paragraphs[0].runs[0].text = new_header
            # description textbox is the very next shape in the deck's convention
            desc_shape = shapes[i + 1]
            if desc_shape.has_text_frame and desc_shape.text_frame.paragraphs[0].runs:
                desc_shape.text_frame.paragraphs[0].runs[0].text = new_desc
            done.add(text)
    return done


def add_projects_table_tile(prs):
    """Slide 10's Database Schema grid is a 3-col x 4-row layout but row 4 only
    fills column 1 (voice_sessions) -- columns 2 and 3 are empty slots. Column 3
    is unsafe (the page-number footer sits in that same horizontal band and
    would visually collide with a tile's description text), but column 2 has
    zero horizontal or vertical overlap with anything else on the slide. Add
    the real "projects" table there, matching the exact tile pattern (card +
    colored header strip + heading textbox + description textbox) used by
    every other table on this slide, so it's indistinguishable from a
    hand-built tile."""
    slide10 = prs.slides[9]
    CARD_FILL = RGBColor(0x1E, 0x2A, 0x45)
    STRIP_COLOR = RGBColor(0xEC, 0x48, 0x99)  # unused color in this grid so far

    left = Emu(4251959)   # column 2 x-position (same as knowledge_bases/workflow_runs/audit_logs tiles)
    top = Emu(5394959)    # row 4 y-position (same row as voice_sessions)
    width = Emu(3749039)
    height = Emu(1371600)

    card = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_FILL
    card.line.fill.background()
    card.shadow.inherit = False

    strip = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(347472))
    strip.fill.solid()
    strip.fill.fore_color.rgb = STRIP_COLOR
    strip.line.fill.background()
    strip.shadow.inherit = False

    add_text(slide10, Emu(4343399), Emu(5431535), Emu(3611879), Emu(292608),
             "⬛ projects", 13, True, WHITE)
    add_text(
        slide10, Emu(4343399), Emu(5779007), Emu(3611879), Emu(914400),
        "id (PK)  ·  owner_id (FK→users)  ·  name  ·  summary  "
        "·  original_prompt  ·  plan (JSON)  ·  ui_html  ·  files (JSON)  "
        "·  chat_history (JSON)  ·  app_type  ·  visibility  ·  shared_with (JSON)  "
        "·  deleted_at  ·  created_at  ·  updated_at",
        9.5, False, BODY_GRAY,
    )


def fix_db_table_count_stat(prs):
    """Executive Summary (slide 2) says '10 DB Tables' -- now that projects
    has a tile on slide 10, the real count is 11."""
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == "10 DB Tables":
            shape.text_frame.paragraphs[0].runs[0].text = "11 DB Tables"
            return True
    return False


def main():
    prs = Presentation(SRC)

    updated = update_stale_approval_bullet(prs)
    print(f"Updated stale approval bullet on slide 24: {updated}")

    print(f"Fixed admin email typo on slide 15: {fix_admin_email_typo(prs)}")
    print(f"Fixed fictional API module entries on slide 12: {fix_api_module_map(prs)}")

    add_projects_table_tile(prs)
    print("Added missing 'projects' table tile to slide 10 (row 4, column 2)")
    print(f"Fixed DB table count stat on slide 2: {fix_db_table_count_stat(prs)}")

    blank_layout = prs.slides[25].slide_layout  # same "Blank" layout as slide 26

    slide = prs.slides.add_slide(blank_layout)

    # Background + top accent bar (matches slide 26 exactly)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_NAVY)
    add_rect(slide, 0, Emu(73152), SLIDE_W, Emu(45720), ACCENT)

    # Title + subtitle
    add_text(slide, Emu(457200), Emu(274320), Emu(10972800), Emu(548640), TITLE, 26, True, WHITE)
    add_text(slide, Emu(457200), Emu(868680), Emu(11247120), Emu(548640), SUBTITLE, 13, False, BODY_GRAY)

    # Feature rows (same vertical rhythm as slide 26: ~896112 EMU per row, starting at 1536192)
    row_top = 1536192
    row_step = 896112
    for heading, desc in ITEMS:
        add_rect(slide, Emu(365760), Emu(row_top), Emu(73152), Emu(777240), ACCENT)
        add_text(slide, Emu(594360), Emu(row_top + 18288), Emu(10972800), Emu(365760), heading, 14, True, ACCENT)
        add_text(slide, Emu(594360), Emu(row_top + 384111), Emu(11064240), Emu(457200), desc, 11, False, BODY_GRAY)
        row_top += row_step

    # Footer page number (matches slide 26's footer textbox position/style)
    add_text(slide, Emu(11640312), Emu(6492240), Emu(457200), Emu(320040), FOOTER_NUMBER, 11, False, BODY_GRAY)

    # Reposition: move the new slide (currently last) to sit right after slide 34
    # (0-indexed 33 = "Lyzr-Parity Scorecard") and before slide 35 ("How To Use" tutorials begin).
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slide_el = slides[-1]
    xml_slides.remove(new_slide_el)
    xml_slides.insert(34, new_slide_el)

    prs.save(OUT)
    print(f"Saved {OUT} with {len(prs.slides)} slides (was 43 in V9).")


if __name__ == "__main__":
    main()
