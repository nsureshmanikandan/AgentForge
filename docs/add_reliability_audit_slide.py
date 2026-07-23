# -*- coding: utf-8 -*-
"""
Adds one more content slide to AgentForge-PresentationV10.pptx documenting the
work done immediately after V10 was built:
  - Deterministic branch input matching (Run modal regenerates its example
    input to match a forced router branch)
  - A page-by-page UI-honesty audit that found 3 independent fake
    tool-integration surfaces and retired the fully non-functional
    Publish-to-Marketplace / Marketplace-Install loop
  - A real production bug found and fixed: /control-plane/stats was crashing
    with a 500 (Decimal + float TypeError), silently making Control Plane and
    Usage & Traceability show 0/blank stat cards
  - Sidebar label consistency pass

Clones the same visual style as the previous "Workflow Builder -- Testing &
Human-in-the-Loop UX" slide (added by add_workflow_testing_slide.py), which
itself matches the deck's native "Latest Platform Updates" slide style: dark
navy background, thin accent top bar, title + description, and N stacked
feature rows.

Inserted right after that slide (now slide 35) and before the "How To Use"
tutorials begin. Saves as AgentForge-PresentationV11.pptx (does not
overwrite V10).
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SRC = r"C:\Users\n.sureshmanikandan\Repo1\AgentForge\docs\AgentForge-PresentationV10.pptx"
OUT = r"C:\Users\n.sureshmanikandan\Repo1\AgentForge\docs\AgentForge-PresentationV11.pptx"

BG_NAVY = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x60, 0x6E, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

SLIDE_W = Emu(12188952)
SLIDE_H = Emu(6858000)

TITLE = "Honest UI Audit & a Real Production Bug — What Changed"
SUBTITLE = (
    "A page-by-page pass asking “is this actually wired to the backend, or does it just look like it "
    "is?” — found three independent fake integration surfaces, a fully non-functional publish loop, "
    "and one genuine crash silently zeroing out two dashboards."
)

ITEMS = [
    (
        "\U0001F3AF  Deterministic Branch Input Matching",
        "The Run modal's example input now regenerates to match whichever router branch is force-selected "
        "(e.g. Critical vs Routine) instead of leaving a stale, narratively mismatched example — a new "
        "target_label parameter on /suggest-input asks the LLM to write an example that specifically "
        "satisfies that router node's own classification criteria.",
    ),
    (
        "\U0001F3F7️  Found: 3 Independent Fake “Connect a Tool” Surfaces",
        "Home's Connect modal, Marketplace's Tool Integrations tab, and Agent Studio's tool checkboxes each "
        "silently did nothing — three separate, unsynchronized implementations of the same idea, none "
        "wired to any real OAuth flow or API credential. All three now carry an honest “not yet active” "
        "label instead of implying real access.",
    ),
    (
        "\U0001F5D1️  Retired the Publish → Marketplace → Install Loop",
        "Agent Studio's “Publish to Marketplace” button and Marketplace's “Published Agents” tab "
        "formed a closed, browser-local loop — both read/wrote the same localStorage key, never touched "
        "the backend, and shared nothing across users or even across browsers on the same machine. Removed "
        "entirely rather than left half-working.",
    ),
    (
        "\U0001F41B  Fixed: Control Plane & Usage Showing Zero",
        "/api/control-plane/stats was silently crashing with a 500 — Postgres returns a float for one "
        "average and a Decimal for another, and adding them raised a TypeError. Both dashboards swallowed "
        "the error and showed 0/blank stat cards despite real data sitting in the tables directly below "
        "them. Cast both averages to float; verified live with real numbers (4 agents, 7 genuine runs).",
    ),
    (
        "\U0001F9F9  Sidebar Naming Consistency",
        "Sidebar labels now match each page's own on-page title exactly instead of drifting from it — "
        "closing a three-way mismatch between the nav label, the page's own heading, and the route/filename "
        "that made it unclear which page you'd actually land on.",
    ),
]

FOOTER_NUMBER = "35a"  # inserted between slide 35 (added previously) and 36 ("How To Use" begins)


def add_rect(slide, left, top, width, height, fill_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size_pt, bold, color_rgb, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color_rgb
    return box


def main():
    prs = Presentation(SRC)
    blank_layout = prs.slides[25].slide_layout  # same "Blank" layout used throughout the deck

    slide = prs.slides.add_slide(blank_layout)

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_NAVY)
    add_rect(slide, 0, Emu(73152), SLIDE_W, Emu(45720), ACCENT)

    add_text(slide, Emu(457200), Emu(274320), Emu(10972800), Emu(548640), TITLE, 24, True, WHITE)
    add_text(slide, Emu(457200), Emu(868680), Emu(11247120), Emu(548640), SUBTITLE, 12, False, BODY_GRAY)

    row_top = 1536192
    row_step = 896112
    for heading, desc in ITEMS:
        add_rect(slide, Emu(365760), Emu(row_top), Emu(73152), Emu(777240), ACCENT)
        add_text(slide, Emu(594360), Emu(row_top + 18288), Emu(10972800), Emu(365760), heading, 14, True, ACCENT)
        add_text(slide, Emu(594360), Emu(row_top + 384111), Emu(11064240), Emu(457200), desc, 10.5, False, BODY_GRAY)
        row_top += row_step

    add_text(slide, Emu(11640312), Emu(6492240), Emu(457200), Emu(320040), FOOTER_NUMBER, 11, False, BODY_GRAY)

    # Move the new slide to sit right after slide 35 (the previous testing/UX slide,
    # 0-indexed 34) and before slide 36 ("How To Use" tutorials begin).
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_slide_el = slides[-1]
    xml_slides.remove(new_slide_el)
    xml_slides.insert(35, new_slide_el)

    prs.save(OUT)
    print(f"Saved {OUT} with {len(prs.slides)} slides (was 44 in V10).")


if __name__ == "__main__":
    main()
