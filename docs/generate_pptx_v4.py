"""
AgentForge — Enterprise Leadership Presentation V4.0
Generates: AgentForge-PresentationV4.0.pptx
Run: python generate_pptx_v4.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Typography ───────────────────────────────────────────────────
FONT_HEADING  = "Segoe UI"
FONT_BODY     = "Segoe UI"
FONT_CODE     = "Consolas"

# ── Colour palette ──────────────────────────────────────────────
C_DARK_BG    = RGBColor(0x08, 0x12, 0x24)
C_BLUE_BG    = RGBColor(0x0b, 0x30, 0x5c)
C_PURPLE_BG  = RGBColor(0x14, 0x0b, 0x30)
C_TEAL_BG    = RGBColor(0x03, 0x28, 0x1e)
C_AMBER_BG   = RGBColor(0x22, 0x10, 0x00)
C_WHITE_BG   = RGBColor(0xf5, 0xf7, 0xfa)
C_SLATE_BG   = RGBColor(0xeb, 0xef, 0xf5)

C_WHITE      = RGBColor(0xff, 0xff, 0xff)
C_OFF_WHITE  = RGBColor(0xd4, 0xdd, 0xee)
C_BLUE_ACC   = RGBColor(0x5b, 0xa3, 0xf5)
C_GREEN_ACC  = RGBColor(0x4a, 0xe0, 0x8c)
C_AMBER_ACC  = RGBColor(0xf9, 0xbb, 0x20)
C_PURPLE_ACC = RGBColor(0xbf, 0xb0, 0xfc)
C_RED_ACC    = RGBColor(0xf5, 0x6e, 0x6e)
C_TEAL_ACC   = RGBColor(0x5e, 0xe2, 0xb4)
C_DARK_TXT   = RGBColor(0x18, 0x24, 0x38)
C_MUTED_TXT  = RGBColor(0x55, 0x66, 0x80)
C_YES        = RGBColor(0x10, 0x98, 0x44)
C_NO         = RGBColor(0xcc, 0x22, 0x22)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────

def add_slide(bg):
    sl = prs.slides.add_slide(blank)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return sl

def rect(sl, x, y, w, h, fill_color, border_color=None, border_w=1.0, radius=None):
    shape = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_w)
    else:
        shape.line.fill.background()
    return shape

def box(sl, x, y, w, h, text, size=14, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, font=FONT_BODY, bg=None):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    return txb

def top_bar(sl, color, height=0.07):
    r = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(height))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()

def slide_header(sl, eyebrow, title, eyebrow_color=None, title_color=C_WHITE, bg=None):
    if bg: rect(sl, 0, 0, 13.33, 7.5, bg)
    ec = eyebrow_color or C_BLUE_ACC
    box(sl, 0.55, 0.3, 12, 0.45, eyebrow.upper(), size=11, bold=True, color=ec, font=FONT_HEADING)
    box(sl, 0.55, 0.72, 12, 0.9, title, size=30, bold=True, color=title_color, font=FONT_HEADING)

def divider(sl, y, color=RGBColor(0x2a, 0x3a, 0x58)):
    r = sl.shapes.add_shape(1, Inches(0.55), Inches(y), Inches(12.23), Inches(0.025))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()

def pill(sl, x, y, w, h, text, bg, fg, size=11):
    r = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    tf = r.text_frame; tf.word_wrap = False
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.name = FONT_BODY; run.font.size = Pt(size)
    run.font.bold = True; run.font.color.rgb = fg

def card_dark(sl, x, y, w, h, title, body, accent=None,
              title_size=15, body_size=12):
    r = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x0e, 0x1e, 0x3c)
    r.line.color.rgb = RGBColor(0x2a, 0x44, 0x72); r.line.width = Pt(0.75)
    if accent:
        a = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(h))
        a.fill.solid(); a.fill.fore_color.rgb = accent; a.line.fill.background()
    box(sl, x+0.14, y+0.13, w-0.22, 0.38, title, size=title_size,
        bold=True, color=C_WHITE, font=FONT_HEADING)
    box(sl, x+0.14, y+0.52, w-0.22, h-0.6, body, size=body_size,
        color=C_OFF_WHITE, font=FONT_BODY)

def card_light(sl, x, y, w, h, title, body, accent_color=None,
               title_size=15, body_size=12):
    r = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = C_WHITE
    r.line.color.rgb = RGBColor(0xc8, 0xd8, 0xec); r.line.width = Pt(0.75)
    if accent_color:
        a = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.06))
        a.fill.solid(); a.fill.fore_color.rgb = accent_color; a.line.fill.background()
    box(sl, x+0.16, y+0.16, w-0.28, 0.40, title, size=title_size,
        bold=True, color=C_DARK_TXT, font=FONT_HEADING)
    box(sl, x+0.16, y+0.56, w-0.28, h-0.65, body, size=body_size,
        color=C_MUTED_TXT, font=FONT_BODY)

def tbl(sl, x, y, w, h, headers, rows, col_widths=None):
    cols = len(headers)
    t    = sl.shapes.add_table(len(rows)+1, cols,
                               Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            t.columns[i].width = Inches(cw)

    def sc(cell, text, bold=False, bg=None, fg=C_DARK_TXT, size=12):
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name  = FONT_BODY
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = fg
        if bg:
            cell.fill.solid(); cell.fill.fore_color.rgb = bg

    for ci, h_txt in enumerate(headers):
        hbg = RGBColor(0x16, 0x52, 0x9a) if ci == 1 else RGBColor(0x18, 0x24, 0x3c)
        sc(t.cell(0, ci), h_txt, bold=True, bg=hbg, fg=C_WHITE, size=12)

    for ri, row in enumerate(rows):
        even_bg = RGBColor(0xee, 0xf4, 0xff)
        for ci, val in enumerate(row):
            bg_c = RGBColor(0xe2, 0xed, 0xff) if ci == 1 else (even_bg if ri%2==0 else C_WHITE)
            fg_c = C_YES if val in ('✅','Yes') else C_NO if val in ('❌','No') else C_DARK_TXT
            sc(t.cell(ri+1, ci), val, bold=(ci==0 or val in ('✅','❌')), bg=bg_c, fg=fg_c, size=12)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_DARK_BG)
top_bar(sl, RGBColor(0x16, 0x52, 0x9a), 0.08)

box(sl, 0.65, 0.5, 7, 0.5, "ENTERPRISE AI AGENT PLATFORM  ·  v4.0",
    size=13, bold=True, color=C_BLUE_ACC, font=FONT_HEADING)
box(sl, 0.65, 0.95, 8, 1.6, "AgentForge",
    size=64, bold=True, color=C_WHITE, font=FONT_HEADING)
box(sl, 0.65, 2.5, 8.5, 0.6,
    "Build · Orchestrate · Govern · Voice-Enable · Plan · Evaluate Production AI Agents",
    size=18, color=C_OFF_WHITE, font=FONT_BODY)
box(sl, 0.65, 3.1, 9, 0.45,
    "Azure-native · Self-hosted · No boilerplate · No vendor lock-in",
    size=14, color=RGBColor(0x70, 0x90, 0xb8), font=FONT_BODY)
divider(sl, 3.65)

pills_data = [
    ("20+ Enterprise Features",   RGBColor(0x12, 0x30, 0x60), C_BLUE_ACC),
    ("Planning Architect",         RGBColor(0x14, 0x06, 0x30), C_PURPLE_ACC),
    ("Voice Agents (Azure TTS)",   RGBColor(0x14, 0x0b, 0x30), C_PURPLE_ACC),
    ("Safety & Guardrails",        RGBColor(0x28, 0x0a, 0x0a), C_RED_ACC),
    ("Evaluations Engine",         RGBColor(0x0a, 0x28, 0x18), C_GREEN_ACC),
]
for i, (txt, bg, fg) in enumerate(pills_data):
    pill(sl, 0.65 + i*2.52, 3.82, 2.38, 0.52, txt, bg, fg, 12)

box(sl, 0.65, 4.55, 12, 0.42,
    "React 18  ·  Vite  ·  TypeScript  ·  ReactFlow  ·  Zustand  ·  TailwindCSS",
    size=13, color=RGBColor(0x50, 0x62, 0x80), font=FONT_BODY)
box(sl, 0.65, 4.95, 12, 0.42,
    "Python 3.12  ·  FastAPI  ·  SQLAlchemy  ·  LangChain  ·  FAISS  ·  PostgreSQL + pgvector  ·  Azure OpenAI  ·  Azure Speech",
    size=13, color=RGBColor(0x50, 0x62, 0x80), font=FONT_BODY)

# watermark
box(sl, 8.2, 0.8, 5.2, 5.0, "AF",
    size=220, bold=True, color=RGBColor(0x0e, 0x1f, 0x3e), font=FONT_HEADING)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive / Leadership Summary
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_DARK_BG)
top_bar(sl, RGBColor(0x16, 0x52, 0x9a))

box(sl, 0.55, 0.28, 5, 0.42, "EXECUTIVE SUMMARY", size=12, bold=True,
    color=C_BLUE_ACC, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.75, "AgentForge — Leadership Brief",
    size=30, bold=True, color=C_WHITE, font=FONT_HEADING)
box(sl, 0.55, 1.42, 12, 0.6,
    "AgentForge is an enterprise-grade, self-hosted AI agent platform enabling teams to design, deploy, "
    "plan, test, voice-enable, and govern intelligent AI automation — visually, without deep ML expertise, "
    "on your Azure subscription.",
    size=14, color=C_OFF_WHITE, font=FONT_BODY)

pts = [
    (C_RED_ACC,    "Problem",
     "Building production AI agents today requires months of engineering, deep Python expertise, "
     "and custom-built safety guardrails — an expensive, high-risk endeavour with no auditability."),
    (C_GREEN_ACC,  "Solution",
     "AgentForge delivers drag-drop visual agent building, Planning Architect (60-second full-stack app generation), "
     "GPT-4o prompt-to-agent generation, built-in Safety & Guardrails, Evaluation engine, Voice Agents via Azure Speech, "
     "and a full observability control plane — out of the box."),
    (C_AMBER_ACC,  "Business Impact",
     "10x faster agent creation  ·  Planning Architect generates full-stack deployable apps in 60 sec  ·  "
     "Zero compliance exposure via always-on PII/safety rules  ·  Voice-enabled agents via Azure Cognitive Services  ·  "
     "Full audit trail for regulated industries  ·  Self-hosted — no SaaS lock-in."),
    (C_PURPLE_ACC, "Competitive Edge",
     "The only platform combining: Planning Architect + Visual Canvas + Azure-native + Self-hosted + Safety Rules + "
     "Voice Agents + Simulation + Evaluations + Workflow Observability. "
     "Beats Flowise, Langflow, Dify, CrewAI, Lyzr, and MS Copilot Studio on all dimensions."),
]
for i, (acc, label, body) in enumerate(pts):
    cy = 2.12 + i*1.22
    r2 = sl.shapes.add_shape(1, Inches(0.55), Inches(cy), Inches(12.23), Inches(1.1))
    r2.fill.solid(); r2.fill.fore_color.rgb = RGBColor(0x08, 0x16, 0x2e)
    r2.line.color.rgb = acc; r2.line.width = Pt(1.0)
    a = sl.shapes.add_shape(1, Inches(0.55), Inches(cy), Inches(0.07), Inches(1.1))
    a.fill.solid(); a.fill.fore_color.rgb = acc; a.line.fill.background()
    box(sl, 0.75, cy+0.08, 2.2, 0.38, label.upper(), size=12, bold=True, color=acc, font=FONT_HEADING)
    box(sl, 0.75, cy+0.46, 11.8, 0.58, body, size=13, color=C_OFF_WHITE, font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Overview (6 cards)
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_WHITE_BG)
box(sl, 0.55, 0.28, 5, 0.42, "PROJECT OVERVIEW", size=12, bold=True,
    color=C_MUTED_TXT, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "What is AgentForge?",
    size=30, bold=True, color=C_DARK_TXT, font=FONT_HEADING)

desc_items = [
    ("Visual Agent Builder",
     "Drag-drop ReactFlow canvas to compose multi-agent pipelines without writing code. "
     "Configure agent name, model, system prompt, tools, and guardrails in a visual sidebar.",
     RGBColor(0x16, 0x52, 0x9a)),
    ("Planning Architect",
     "Describe your app in natural language — GPT-4o generates a full project plan, working React UI preview, "
     "and a deployable full-stack ZIP (React+FastAPI+FAISS+PostgreSQL+Docker) in one session.",
     RGBColor(0x5a, 0x1a, 0xa0)),
    ("Multi-Agent Orchestration",
     "Manager agents delegate to worker agents. Each worker runs its own guardrail pass. "
     "Supports sequential, parallel, and hierarchical pipeline topologies.",
     RGBColor(0x0c, 0x5e, 0x48)),
    ("Safety & Evaluations",
     "8 configurable safety rules (PII, hallucination, toxicity, prompt injection). "
     "Evaluation engine runs automated test suites and scores agents before promotion to production.",
     RGBColor(0x8e, 0x28, 0x28)),
    ("Voice Agents",
     "Azure Cognitive Services Speech SDK integration: text-to-speech synthesis with 8 Neural voices "
     "(Jenny, Sonia, Natasha, Neerja and more). Full telephony config per agent.",
     RGBColor(0x44, 0x18, 0x9a)),
    ("Enterprise Control Plane",
     "Live stats dashboard, workflow observability traces, agent versioning, full audit trail, "
     "RBAC (Admin/Developer/Viewer), OpenTelemetry tracing to Jaeger/Azure Monitor/GCP/Datadog.",
     RGBColor(0x16, 0x52, 0x9a)),
]
for i, (title, body, acc) in enumerate(desc_items):
    col = i % 3; row = i // 3
    card_light(sl, 0.55 + col*4.27, 1.62 + row*2.72, 4.0, 2.55,
               title, body, accent_color=acc, title_size=15, body_size=12)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — The Problem We Solve
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_PURPLE_BG)
top_bar(sl, RGBColor(0x6c, 0x28, 0xe0))
slide_header(sl, "The Challenge", "Enterprise AI agents are broken by design")

probs = [
    ("Code-Heavy Barrier",
     "Most frameworks demand deep Python & ML expertise "
     "before a single agent reaches production. "
     "Average time-to-first-agent: 8 weeks.",
     RGBColor(0x6c, 0x28, 0xe0)),
    ("Zero Compliance Controls",
     "Open-source tools ship with no PII protection, "
     "no hallucination checks, and no toxicity filters — "
     "a critical compliance risk in regulated industries.",
     RGBColor(0x6c, 0x28, 0xe0)),
    ("No Observability",
     "No way to observe, version, evaluate, test, or audit "
     "what agents actually do after they are deployed. "
     "Production incidents are invisible until users complain.",
     RGBColor(0x6c, 0x28, 0xe0)),
]
for i, (t, b, acc) in enumerate(probs):
    card_dark(sl, 0.55 + i*4.27, 2.05, 3.98, 3.7, t, b,
              accent=acc, title_size=16, body_size=13)

r = sl.shapes.add_shape(1, Inches(0.55), Inches(6.0), Inches(12.23), Inches(1.1))
r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x10, 0x06, 0x28)
r.line.color.rgb = RGBColor(0x6c, 0x28, 0xe0); r.line.width = Pt(1.5)
box(sl, 0.75, 6.08, 12, 0.42, "AgentForge solves all three — and more",
    size=16, bold=True, color=C_WHITE, font=FONT_HEADING)
box(sl, 0.75, 6.5, 12, 0.42,
    "Planning Architect  ·  Visual builder  ·  8 Safety rules  ·  Evaluations engine  ·  Voice agents  ·  Workflow Observability",
    size=13, color=C_OFF_WHITE, font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — ROI & Business Value
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_WHITE_BG)
box(sl, 0.55, 0.28, 6, 0.42, "BUSINESS VALUE & ROI",
    size=12, bold=True, color=C_MUTED_TXT, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "Return on Investment — Why AgentForge Pays for Itself",
    size=28, bold=True, color=C_DARK_TXT, font=FONT_HEADING)

# Big numbers
big_nums = [
    ("10x",  "Faster Agent\nCreation",        RGBColor(0x16, 0x52, 0x9a), C_BLUE_ACC),
    ("60s",  "Full-Stack App\nGeneration",    RGBColor(0x3a, 0x08, 0x6a), C_PURPLE_ACC),
    ("100%", "Audit Coverage\nPer Run",        RGBColor(0x4a, 0x30, 0xb0), C_PURPLE_ACC),
    ("8",    "Safety Rules\nOut of the Box",  RGBColor(0x7a, 0x20, 0x20), C_RED_ACC),
]
for i, (num, label, bg, acc) in enumerate(big_nums):
    rx = 0.55 + i*3.2
    r2 = sl.shapes.add_shape(1, Inches(rx), Inches(1.62), Inches(3.0), Inches(1.42))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg
    r2.line.color.rgb = acc; r2.line.width = Pt(1.0)
    box(sl, rx+0.1, 1.66, 2.8, 0.78, num, size=46, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_HEADING)
    box(sl, rx+0.1, 2.42, 2.8, 0.56, label, size=12,
        color=acc, align=PP_ALIGN.CENTER, font=FONT_BODY)

roi_cards = [
    ("Development Speed",
     "Traditional builds: 6-12 weeks of engineering.\n"
     "AgentForge: prompt-to-agent in under 60 seconds.\n"
     "Estimated saving: 3-6 developer-weeks per agent.",
     RGBColor(0x16, 0x52, 0x9a)),
    ("Planning Architect Value",
     "60-sec full-stack app generation: React+FastAPI+FAISS+Docker.\n"
     "KB-grounded RAG with two-pass extraction — no hallucination.\n"
     "Deployable ZIP with one docker-compose command.",
     RGBColor(0x5a, 0x18, 0xa0)),
    ("SaaS Cost Avoidance",
     "Lyzr: ~$2,000+/month SaaS. Dify Cloud: per-seat.\n"
     "AgentForge: self-hosted on your Azure subscription.\n"
     "Break-even at first agent deployed.",
     RGBColor(0x78, 0x46, 0x08)),
    ("Compliance Risk Reduction",
     "8 always-on safety rules (PII, hallucination, toxicity,\n"
     "prompt injection) eliminate GDPR/HIPAA exposure.\n"
     "Audit logs satisfy regulatory review in minutes.",
     RGBColor(0x0c, 0x5e, 0x48)),
    ("Quality Assurance",
     "Evaluation engine scores agents pre-deployment.\n"
     "Only agents meeting pass-rate threshold promote.\n"
     "Workflow Observability surfaces failures instantly.",
     RGBColor(0x7a, 0x20, 0x20)),
    ("Talent Leverage",
     "Non-ML staff build production agents visually.\n"
     "Reduces dependency on scarce AI engineers.\n"
     "Config-driven safety — zero ML knowledge needed.",
     RGBColor(0x16, 0x52, 0x9a)),
]
for i, (title, body, acc) in enumerate(roi_cards):
    col = i % 3; row = i // 3
    card_light(sl, 0.55 + col*4.27, 3.2 + row*2.05, 4.0, 1.95,
               title, body, accent_color=acc, title_size=14, body_size=12)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture Diagram
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_DARK_BG)
top_bar(sl, RGBColor(0x16, 0x52, 0x9a))
box(sl, 0.55, 0.2, 6, 0.42, "TECHNICAL ARCHITECTURE", size=12, bold=True,
    color=C_BLUE_ACC, font=FONT_HEADING)
box(sl, 0.55, 0.6, 12, 0.75, "AgentForge — Layered Architecture",
    size=30, bold=True, color=C_WHITE, font=FONT_HEADING)

layers = [
    (6.3,  0.65, RGBColor(0x04, 0x18, 0x34), RGBColor(0x16, 0x52, 0x9a),
     "AZURE AI LAYER",
     "Azure OpenAI GPT-4o  ·  Azure OpenAI GPT-4.5  ·  Azure AI Search  ·  Azure Cognitive Services Speech (TTS/STT)"),
    (5.58, 0.65, RGBColor(0x02, 0x1c, 0x14), RGBColor(0x0c, 0x5e, 0x48),
     "DATA LAYER",
     "PostgreSQL 16 + pgvector  ·  FAISS (Planning Architect RAG)  ·  Alembic Migrations  ·  SQLAlchemy 2.0 Async ORM"),
    (4.14, 1.38, RGBColor(0x10, 0x06, 0x2c), RGBColor(0x4a, 0x30, 0xb0),
     "CORE SERVICES",
     "AgentOrchestrator  ·  PlanningArchitect (GPT-4o plan + UI gen + ZIP export)  ·  GuardrailsEngine (Presidio)  ·  "
     "RAGEngine  ·  SimulationRunner  ·  ToolRegistry  ·  PromptToAgent  ·  SafetyRules (8)  ·  EvaluationEngine  ·  VoiceService"),
    (2.98, 1.1, RGBColor(0x22, 0x0e, 0x00), RGBColor(0x78, 0x46, 0x08),
     "API LAYER — FastAPI + JWT RBAC",
     "Agents  ·  Auth  ·  RAG  ·  Simulation  ·  Tools  ·  Control Plane  ·  Safety  ·  Evaluations  ·  Voice  ·  Architect  ·  WorkflowRuns"),
    (1.86, 1.06, RGBColor(0x06, 0x12, 0x26), RGBColor(0x5b, 0xa3, 0xf5),
     "FRONTEND — React 18 + Vite + TypeScript",
     "ReactFlow Canvas  ·  Agent Studio  ·  Planning Architect  ·  Workflow Observability  ·  Safety & Guardrails  ·  Evaluations  ·  Voice Agents  ·  Dashboard"),
]
for (ly, lh, bg, border, label, items) in layers:
    r2 = sl.shapes.add_shape(1, Inches(0.55), Inches(ly), Inches(10.5), Inches(lh))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg
    r2.line.color.rgb = border; r2.line.width = Pt(1.2)
    box(sl, 0.72, ly+0.05, 3.5, 0.32, label, size=11, bold=True, color=border, font=FONT_HEADING)
    box(sl, 0.72, ly+0.36, 10.2, lh-0.42, items, size=12, color=C_OFF_WHITE, font=FONT_BODY)

# Observability sidebar
rect(sl, 11.22, 1.86, 1.66, 5.14, RGBColor(0x02, 0x14, 0x0c),
     border_color=RGBColor(0x18, 0x7a, 0x5e))
box(sl, 11.28, 1.92, 1.52, 0.32, "OBSERVABILITY",
    size=9, bold=True, color=C_TEAL_ACC, font=FONT_HEADING)
for i, item in enumerate(["OpenTelemetry", "Jaeger", "Azure Monitor", "GCP Trace", "Datadog APM", "Workflow Runs"]):
    box(sl, 11.28, 2.32 + i*0.62, 1.48, 0.56, item, size=11, color=C_OFF_WHITE, font=FONT_BODY)

# Docker row
rect(sl, 0.55, 7.12, 10.5, 0.32, RGBColor(0x02, 0x10, 0x08),
     border_color=RGBColor(0x18, 0x7a, 0x5e))
box(sl, 0.72, 7.14, 10.2, 0.28,
    "Docker Compose  —  postgres + pgvector + backend + frontend + jaeger",
    size=11, color=C_TEAL_ACC, font=FONT_CODE)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Planning Architect Deep Dive
# ════════════════════════════════════════════════════════════════
sl = add_slide(RGBColor(0x10, 0x06, 0x28))
top_bar(sl, C_PURPLE_ACC, 0.08)
box(sl, 0.55, 0.28, 8, 0.42, "NEW FEATURE SPOTLIGHT — v4.0", size=12, bold=True,
    color=C_PURPLE_ACC, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "Planning Architect — 60-Second Full-Stack App Generation",
    size=28, bold=True, color=C_WHITE, font=FONT_HEADING)

# Stats row
arch_stats = [
    ("60s",  "Describe → Full\nProject Plan",    RGBColor(0x22, 0x08, 0x44), C_PURPLE_ACC),
    ("1",    "Working React UI\nPreview",          RGBColor(0x08, 0x22, 0x44), C_BLUE_ACC),
    ("2",    "KB Extraction\nPasses",              RGBColor(0x08, 0x28, 0x1c), C_GREEN_ACC),
    ("1 ZIP","Full-Stack Deploy\nReady",           RGBColor(0x22, 0x18, 0x00), C_AMBER_ACC),
]
for i, (num, label, bg, acc) in enumerate(arch_stats):
    rx = 0.55 + i*3.2
    r2 = sl.shapes.add_shape(1, Inches(rx), Inches(1.55), Inches(3.0), Inches(1.2))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg
    r2.line.color.rgb = acc; r2.line.width = Pt(1.0)
    box(sl, rx+0.1, 1.6, 2.8, 0.62, num, size=36, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_HEADING)
    box(sl, rx+0.1, 2.2, 2.8, 0.48, label, size=11,
        color=acc, align=PP_ALIGN.CENTER, font=FONT_BODY)

# Left — How it works
box(sl, 0.55, 2.92, 5.8, 0.4, "HOW IT WORKS — 5 STEPS", size=12, bold=True,
    color=C_PURPLE_ACC, font=FONT_HEADING)
arch_steps = [
    ("1", "Describe Your App",      "Type what you want to build in plain English. Attach PDFs/DOCXs for KB-grounded generation."),
    ("2", "AI Generates Full Plan", "GPT-4o produces: project summary, phases, features, tech stack, architecture diagram in JSON."),
    ("3", "Visual UI Preview",      "A working React chatbot/dashboard renders as a sandboxed HTML preview — live in browser."),
    ("4", "Two-Pass KB Extraction", "Pass 1: extract FAQ_DATA/topics as structured JSON. Pass 2: inject into React template — no hallucination."),
    ("5", "Download Full-Stack ZIP","Get React+Vite+TS frontend, FastAPI backend, FAISS RAG, PostgreSQL schema, Docker Compose — deploy anywhere."),
]
for i, (num, title, body) in enumerate(arch_steps):
    cy = 3.4 + i*0.82
    r2 = sl.shapes.add_shape(1, Inches(0.55), Inches(cy), Inches(6.0), Inches(0.73))
    r2.fill.solid(); r2.fill.fore_color.rgb = RGBColor(0x18, 0x0a, 0x38)
    r2.line.color.rgb = C_PURPLE_ACC; r2.line.width = Pt(0.5)
    pill(sl, 0.62, cy+0.12, 0.44, 0.44, num, C_PURPLE_ACC, C_WHITE, size=13)
    box(sl, 1.18, cy+0.04, 5.2, 0.3, title, size=13, bold=True, color=C_WHITE, font=FONT_HEADING)
    box(sl, 1.18, cy+0.37, 5.2, 0.32, body, size=11, color=C_OFF_WHITE, font=FONT_BODY)

# Right — ZIP contents
box(sl, 7.0, 2.92, 5.8, 0.4, "FULL-STACK ZIP CONTENTS", size=12, bold=True,
    color=C_PURPLE_ACC, font=FONT_HEADING)
zip_items = [
    ("src/App.tsx",                  "React + TypeScript chatbot UI", C_BLUE_ACC),
    ("backend/main.py",              "FastAPI application entry point", C_GREEN_ACC),
    ("backend/app/rag.py",           "FAISS + Azure OpenAI RAG engine", C_GREEN_ACC),
    ("backend/app/models.py",        "PostgreSQL tables (auto-created on startup)", C_GREEN_ACC),
    ("backend/requirements.txt",     "Python dependencies pinned", C_TEAL_ACC),
    ("backend/.env.example",         "Azure keys + DB URL template", C_AMBER_ACC),
    ("docker-compose.yml",           "One-command deploy: postgres + backend + frontend", C_AMBER_ACC),
    ("sandbox.html",                 "Standalone demo — no install needed", C_PURPLE_ACC),
]
for i, (fname, desc, acc) in enumerate(zip_items):
    cy = 3.4 + i*0.52
    r2 = sl.shapes.add_shape(1, Inches(7.0), Inches(cy), Inches(6.1), Inches(0.44))
    r2.fill.solid(); r2.fill.fore_color.rgb = RGBColor(0x14, 0x08, 0x32)
    r2.line.color.rgb = acc; r2.line.width = Pt(0.4)
    box(sl, 7.1, cy+0.02, 2.6, 0.3, fname, size=11, bold=True, color=acc, font=FONT_CODE)
    box(sl, 9.78, cy+0.04, 3.2, 0.3, desc, size=11, color=C_OFF_WHITE, font=FONT_BODY)

# Session management note
rect(sl, 0.55, 7.12, 12.23, 0.32, RGBColor(0x18, 0x0a, 0x38),
     border_color=C_PURPLE_ACC)
box(sl, 0.72, 7.14, 12, 0.26,
    "Session Management: Multiple named sessions (#1 · New Session · Jul 12)  ·  Scrollable list  ·  Delete per session  ·  Document upload per session",
    size=12, color=C_PURPLE_ACC, font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Full-Stack Download (ZIP tree)
# ════════════════════════════════════════════════════════════════
sl = add_slide(RGBColor(0x02, 0x1c, 0x18))
top_bar(sl, C_TEAL_ACC, 0.08)
box(sl, 0.55, 0.28, 8, 0.42, "FEATURE SPOTLIGHT — FULL-STACK DOWNLOAD", size=12, bold=True,
    color=C_TEAL_ACC, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "Deploy-Ready ZIP: React + FastAPI + FAISS + PostgreSQL + Docker",
    size=26, bold=True, color=C_WHITE, font=FONT_HEADING)

# Left — Frontend files
rect(sl, 0.55, 1.55, 5.7, 5.35, RGBColor(0x04, 0x22, 0x1a),
     border_color=RGBColor(0x18, 0x9a, 0x74))
box(sl, 0.72, 1.62, 5.4, 0.38, "FRONTEND  (React 18 + Vite + TypeScript)",
    size=12, bold=True, color=C_TEAL_ACC, font=FONT_HEADING)
fe_files = [
    "my-app/",
    "  src/",
    "    App.tsx          # Main chatbot / dashboard UI",
    "    components/      # Chat, MessageList, InputBar",
    "    api.ts           # Axios client → backend",
    "  public/",
    "  index.html",
    "  vite.config.ts",
    "  tsconfig.json",
    "  package.json       # React 18 + TypeScript + Vite",
    "sandbox.html         # Standalone no-install demo",
]
for i, line in enumerate(fe_files):
    box(sl, 0.72, 2.08 + i*0.44, 5.4, 0.4, line, size=12,
        color=C_TEAL_ACC if line.endswith("/") else C_OFF_WHITE,
        font=FONT_CODE, bold=line.endswith("/"))

# Right — Backend files
rect(sl, 6.6, 1.55, 6.18, 5.35, RGBColor(0x04, 0x1a, 0x22),
     border_color=RGBColor(0x18, 0x74, 0x9a))
box(sl, 6.78, 1.62, 5.8, 0.38, "BACKEND  (FastAPI + FAISS + PostgreSQL)",
    size=12, bold=True, color=C_BLUE_ACC, font=FONT_HEADING)
be_files = [
    "backend/",
    "  main.py            # FastAPI entry — routers, CORS",
    "  app/",
    "    rag.py           # FAISS + Azure OpenAI RAG",
    "    models.py        # PostgreSQL tables (auto-init)",
    "    routes/          # /ask, /upload, /health",
    "  requirements.txt   # fastapi, langchain, faiss-cpu…",
    "  .env.example       # Azure keys + DB URL template",
    "  Dockerfile",
    "docker-compose.yml   # postgres + backend + frontend",
    "  → docker-compose up --build",
]
for i, line in enumerate(be_files):
    box(sl, 6.78, 2.08 + i*0.44, 5.8, 0.4, line, size=12,
        color=C_BLUE_ACC if line.endswith("/") else C_OFF_WHITE,
        font=FONT_CODE, bold=line.endswith("/"))

# Bottom deploy banner
rect(sl, 0.55, 7.0, 12.23, 0.44, RGBColor(0x01, 0x12, 0x10),
     border_color=C_TEAL_ACC)
box(sl, 0.72, 7.04, 12, 0.34,
    "ONE COMMAND DEPLOY:   docker-compose up --build   →   Frontend: :5173   Backend: :8000   PostgreSQL: :5432",
    size=13, bold=True, color=C_TEAL_ACC, font=FONT_CODE)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Safety & Guardrails
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_DARK_BG)
top_bar(sl, C_RED_ACC, 0.08)
box(sl, 0.55, 0.28, 8, 0.42, "FEATURE SPOTLIGHT", size=12, bold=True,
    color=C_RED_ACC, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "Safety & Guardrails Engine",
    size=30, bold=True, color=C_WHITE, font=FONT_HEADING)
box(sl, 0.55, 1.44, 12, 0.5,
    "8 configurable safety rules across 4 categories — every agent run is scanned before output is returned",
    size=14, color=C_OFF_WHITE, font=FONT_BODY)

# Stats row
stats = [
    ("1,247", "Total Requests\nMonitored",   RGBColor(0x08, 0x20, 0x44), C_BLUE_ACC),
    ("23",    "Blocked\nRequests",           RGBColor(0x28, 0x08, 0x08), C_RED_ACC),
    ("89",    "PII\nRedactions",             RGBColor(0x22, 0x10, 0x3c), C_PURPLE_ACC),
    ("14",    "Hallucination\nFlags",        RGBColor(0x22, 0x18, 0x00), C_AMBER_ACC),
]
for i, (num, label, bg, acc) in enumerate(stats):
    rx = 0.55 + i*3.2
    r2 = sl.shapes.add_shape(1, Inches(rx), Inches(2.1), Inches(3.0), Inches(1.2))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg
    r2.line.color.rgb = acc; r2.line.width = Pt(1.0)
    box(sl, rx+0.1, 2.15, 2.8, 0.65, num, size=40, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_HEADING)
    box(sl, rx+0.1, 2.78, 2.8, 0.46, label, size=12,
        color=acc, align=PP_ALIGN.CENTER, font=FONT_BODY)

# Rules grid
rules = [
    ("PII Detection",      "Critical · Blocks & redacts SSN, email,\nphone, credit card via Microsoft Presidio.",  C_RED_ACC),
    ("Hallucination Check","High · Flags uncertain language patterns\nbefore output reaches the end user.",        C_AMBER_ACC),
    ("Toxicity Filter",    "High · Screens for harmful, offensive or\ndegrading content in every response.",       C_RED_ACC),
    ("Prompt Injection",   "Critical · Blocks jailbreak & instruction\noverride attempts in user inputs.",         C_RED_ACC),
    ("Data Leakage Guard", "High · Prevents confidential data from\nbeing exposed in agent responses.",            C_AMBER_ACC),
    ("Off-Topic Filter",   "Medium · Keeps agents focused on their\nconfigured domain and persona.",               C_BLUE_ACC),
    ("Rate Limit Guard",   "Medium · Prevents abuse and runaway\ncost from excessive agent calls.",                C_BLUE_ACC),
    ("Output Length",      "Low · Caps responses to avoid token\nwaste and truncation issues.",                    C_TEAL_ACC),
]
for i, (name, body, acc) in enumerate(rules):
    col = i % 4; row = i // 4
    card_dark(sl, 0.55 + col*3.2, 3.5 + row*1.72, 3.0, 1.62,
              name, body, accent=acc, title_size=14, body_size=12)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Evaluations Engine
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_WHITE_BG)
box(sl, 0.55, 0.28, 6, 0.42, "FEATURE SPOTLIGHT", size=12, bold=True,
    color=C_MUTED_TXT, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "Evaluations Engine — Automated Agent QA",
    size=28, bold=True, color=C_DARK_TXT, font=FONT_HEADING)

# Left — How it works
box(sl, 0.55, 1.55, 5.8, 0.42, "HOW IT WORKS", size=12, bold=True,
    color=C_MUTED_TXT, font=FONT_HEADING)
steps_eval = [
    ("1", "Define Test Cases",    "JSON array of {input, expected} pairs targeting your agent's use cases"),
    ("2", "Run Evaluation",       "POST /api/evaluations/runs  ->  engine scores each test case pass/fail"),
    ("3", "Score & Badge",        ">=80% = green (production-ready)  ·  60-79% = yellow  ·  <60% = red (blocked)"),
    ("4", "Review History",       "Full run history with agent ID, timestamp, passed/total, and per-case results"),
]
for i, (num, title, body) in enumerate(steps_eval):
    cy = 2.05 + i*1.18
    r2 = sl.shapes.add_shape(1, Inches(0.55), Inches(cy), Inches(5.8), Inches(1.08))
    r2.fill.solid(); r2.fill.fore_color.rgb = RGBColor(0xf0, 0xf5, 0xff)
    r2.line.color.rgb = RGBColor(0x16, 0x52, 0x9a); r2.line.width = Pt(0.75)
    pill(sl, 0.62, cy+0.3, 0.48, 0.48, num, RGBColor(0x16, 0x52, 0x9a), C_WHITE, size=14)
    box(sl, 1.22, cy+0.1, 5.0, 0.36, title, size=14, bold=True, color=C_DARK_TXT, font=FONT_HEADING)
    box(sl, 1.22, cy+0.5, 4.9, 0.5, body, size=12, color=C_MUTED_TXT, font=FONT_BODY)

# Right — Benefits
box(sl, 6.7, 1.55, 6.1, 0.42, "BUSINESS BENEFITS", size=12, bold=True,
    color=C_MUTED_TXT, font=FONT_HEADING)
benefits = [
    ("Pre-Production Gate",
     "Agents must pass score threshold before\nbeing promoted to production. Prevents\ndefective agents from reaching users.",
     RGBColor(0x16, 0x52, 0x9a)),
    ("Regression Testing",
     "Re-run eval suite after every agent update.\nCatch regressions before users do.\nFull pass/fail diff in the history table.",
     RGBColor(0x0c, 0x5e, 0x48)),
    ("Compliance Evidence",
     "Each eval run is timestamped and persisted.\nProvides auditable evidence that agents\nmeet quality standards for regulated industries.",
     RGBColor(0x4a, 0x30, 0xb0)),
    ("No Custom Test Framework",
     "JSON test cases. No pytest. No ML pipeline.\nProduct teams write tests — not engineers.\nZero setup beyond defining inputs.",
     RGBColor(0x78, 0x46, 0x08)),
]
for i, (t, b, acc) in enumerate(benefits):
    col = i % 2; row = i // 2
    card_light(sl, 6.7 + col*3.17, 2.05 + row*2.5, 2.98, 2.35,
               t, b, accent_color=acc, title_size=14, body_size=12)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Voice Agents
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_PURPLE_BG)
top_bar(sl, C_PURPLE_ACC, 0.08)
box(sl, 0.55, 0.28, 8, 0.42, "FEATURE SPOTLIGHT", size=12, bold=True,
    color=C_PURPLE_ACC, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "Voice Agents — Azure Cognitive Services Speech",
    size=30, bold=True, color=C_WHITE, font=FONT_HEADING)
box(sl, 0.55, 1.44, 11, 0.5,
    "Real-time TTS synthesis via Azure Neural voices  ·  Per-agent telephony configuration  ·  STT ready",
    size=14, color=C_OFF_WHITE, font=FONT_BODY)

# Left — Voices grid
box(sl, 0.55, 2.0, 6, 0.4, "8 AZURE NEURAL VOICES INCLUDED", size=12, bold=True,
    color=C_PURPLE_ACC, font=FONT_HEADING)
voices = [
    ("Jenny (US)",     "Female", C_PURPLE_ACC),
    ("Guy (US)",       "Male",   C_BLUE_ACC),
    ("Sonia (UK)",     "Female", C_PURPLE_ACC),
    ("Natasha (AU)",   "Female", C_PURPLE_ACC),
    ("Neerja (India)", "Female", C_PURPLE_ACC),
    ("Elvira (ES)",    "Female", C_TEAL_ACC),
    ("Denise (FR)",    "Female", C_TEAL_ACC),
    ("Katja (DE)",     "Female", C_TEAL_ACC),
]
for i, (name, gender, acc) in enumerate(voices):
    col = i % 4; row = i // 4
    r2 = sl.shapes.add_shape(1, Inches(0.55 + col*1.62), Inches(2.52 + row*0.82), Inches(1.55), Inches(0.72))
    r2.fill.solid(); r2.fill.fore_color.rgb = RGBColor(0x18, 0x0c, 0x3a)
    r2.line.color.rgb = acc; r2.line.width = Pt(0.75)
    box(sl, 0.62 + col*1.62, 2.56 + row*0.82, 1.4, 0.3, name, size=11, bold=True, color=C_WHITE, font=FONT_BODY)
    box(sl, 0.62 + col*1.62, 2.84 + row*0.82, 1.4, 0.24, gender, size=10, color=acc, font=FONT_BODY)

# Left — Config
box(sl, 0.55, 4.3, 6, 0.4, "PER-AGENT VOICE CONFIG", size=12, bold=True,
    color=C_PURPLE_ACC, font=FONT_HEADING)
cfg_items = ["TTS Engine (Azure / OpenAI)", "STT Engine (Azure / Whisper)",
             "TTS Voice selection", "Silence Timeout (ms)", "Max Call Duration (s)"]
for i, item in enumerate(cfg_items):
    box(sl, 0.75, 4.76 + i*0.46, 5.6, 0.4, f"•  {item}", size=13, color=C_OFF_WHITE, font=FONT_BODY)

# Right — How it works
box(sl, 7.2, 2.0, 5.8, 0.4, "HOW IT WORKS", size=12, bold=True,
    color=C_PURPLE_ACC, font=FONT_HEADING)
how_steps = [
    ("Select Voice",         "Choose from 8 Azure Neural voices on the Voice Agents page"),
    ("Configure Agent",      "Set TTS/STT engine, silence timeout, max duration, enable voice"),
    ("Test Preview",         "Type sample text -> click Play Sample -> Azure Speech SDK synthesises MP3"),
    ("Save Config",          "PUT /api/voice/configs/{agent_id}  stores settings per agent"),
    ("Deploy",               "Agent answers inbound calls or places outbound calls via Azure telephony"),
]
for i, (label, body) in enumerate(how_steps):
    cy = 2.5 + i*0.96
    r2 = sl.shapes.add_shape(1, Inches(7.2), Inches(cy), Inches(5.8), Inches(0.86))
    r2.fill.solid(); r2.fill.fore_color.rgb = RGBColor(0x14, 0x08, 0x2e)
    r2.line.color.rgb = C_PURPLE_ACC; r2.line.width = Pt(0.5)
    box(sl, 7.36, cy+0.06, 1.8, 0.3, label, size=12, bold=True, color=C_PURPLE_ACC, font=FONT_HEADING)
    box(sl, 7.36, cy+0.4, 5.5, 0.38, body, size=12, color=C_OFF_WHITE, font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Agent Studio & Type Filters
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_BLUE_BG)
top_bar(sl, C_BLUE_ACC, 0.08)
slide_header(sl, "Feature Spotlight", "Agent Studio — Visual Builder & Type Filters")

# Left features
feats = [
    ("Agent Type Filters",
     "Filter tabs: All · Agent · Managerial · Superflow\n"
     "Each agent has agent_type stored in PostgreSQL.\n"
     "Counts update live as agents are created.",
     C_BLUE_ACC),
    ("Drag-Drop Canvas",
     "ReactFlow v12 with animated edges, dark mode,\n"
     "and dot-grid background. No code required to build\n"
     "complex multi-agent pipeline topologies.",
     C_BLUE_ACC),
    ("Prompt-to-Agent AI",
     "Describe your agent in plain English — GPT-4o\n"
     "generates name, system prompt, tools & guardrails\n"
     "in under 60 seconds.",
     C_GREEN_ACC),
    ("Manager / Worker Pattern",
     "Manager agent delegates tasks to worker agents.\n"
     "Each worker runs its own guardrail pass before\n"
     "results are aggregated by the manager.",
     C_AMBER_ACC),
]
for i, (t, b, acc) in enumerate(feats):
    card_dark(sl, 0.55, 2.0 + i*1.32, 5.7, 1.22, t, b,
              accent=acc, title_size=15, body_size=12)

# Right — canvas topology diagram
rect(sl, 6.6, 2.0, 6.25, 5.18, RGBColor(0x04, 0x14, 0x30),
     border_color=RGBColor(0x2a, 0x5a, 0xa0))
box(sl, 6.78, 2.1, 6, 0.36, "AGENT STUDIO CANVAS LAYOUT",
    size=11, bold=True, color=RGBColor(0x50, 0x88, 0xcc), font=FONT_HEADING)

# Filter tabs row
tabs = [("All 12", RGBColor(0x16, 0x52, 0x9a)), ("Agent 10", RGBColor(0x0c, 0x44, 0x7a)),
        ("Managerial 1", RGBColor(0x0c, 0x44, 0x7a)), ("Superflow 1", RGBColor(0x0c, 0x44, 0x7a))]
for i, (label, bg) in enumerate(tabs):
    r2 = sl.shapes.add_shape(1, Inches(6.78 + i*1.5), Inches(2.55), Inches(1.42), Inches(0.38))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg; r2.line.fill.background()
    box(sl, 6.82 + i*1.5, 2.58, 1.35, 0.3, label, size=11, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_BODY)

nodes_data = [
    (7.0,  3.25, 1.7, 0.52, "Input Node",    RGBColor(0x0e, 0x22, 0x4a)),
    (8.88, 3.25, 1.7, 0.52, "Manager",       RGBColor(0x16, 0x52, 0x9a)),
    (10.7, 3.25, 1.7, 0.52, "Output",        RGBColor(0x0e, 0x22, 0x4a)),
    (7.5,  4.3,  1.8, 0.52, "Worker A",      RGBColor(0x0c, 0x44, 0x28)),
    (9.5,  4.3,  1.8, 0.52, "Worker B",      RGBColor(0x0c, 0x44, 0x28)),
    (7.5,  5.35, 1.8, 0.52, "Worker C",      RGBColor(0x32, 0x18, 0x00)),
    (9.5,  5.35, 1.8, 0.52, "Tool: Slack",   RGBColor(0x32, 0x18, 0x00)),
]
for nx, ny, nw, nh, nt, nc in nodes_data:
    r2 = sl.shapes.add_shape(1, Inches(nx), Inches(ny), Inches(nw), Inches(nh))
    r2.fill.solid(); r2.fill.fore_color.rgb = nc
    r2.line.color.rgb = RGBColor(0x3a, 0x6a, 0xcc); r2.line.width = Pt(0.75)
    tf = r2.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run(); run.text = nt
    run.font.name = FONT_BODY; run.font.size = Pt(13)
    run.font.bold = True; run.font.color.rgb = C_WHITE

box(sl, 6.75, 6.5, 6, 0.42,
    "Unlimited nodes per workflow  ·  Single or multi-agent modes",
    size=12, color=RGBColor(0x70, 0xa0, 0xdd), font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Workflow Observability
# ════════════════════════════════════════════════════════════════
sl = add_slide(RGBColor(0x02, 0x0e, 0x28))
top_bar(sl, C_BLUE_ACC, 0.08)
box(sl, 0.55, 0.28, 8, 0.42, "NEW FEATURE SPOTLIGHT — v4.0", size=12, bold=True,
    color=C_BLUE_ACC, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "Workflow Observability — Full Execution Trace Dashboard",
    size=26, bold=True, color=C_WHITE, font=FONT_HEADING)

# Stats bar
obs_stats = [
    ("1,482", "Total Runs",       RGBColor(0x04, 0x18, 0x34), C_BLUE_ACC),
    ("1,401", "Completed",        RGBColor(0x04, 0x22, 0x14), C_GREEN_ACC),
    ("81",    "Failed",           RGBColor(0x28, 0x06, 0x06), C_RED_ACC),
    ("2.3s",  "Avg Duration",     RGBColor(0x22, 0x18, 0x00), C_AMBER_ACC),
]
for i, (num, label, bg, acc) in enumerate(obs_stats):
    rx = 0.55 + i*3.2
    r2 = sl.shapes.add_shape(1, Inches(rx), Inches(1.55), Inches(3.0), Inches(1.15))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg
    r2.line.color.rgb = acc; r2.line.width = Pt(1.0)
    box(sl, rx+0.1, 1.6, 2.8, 0.6, num, size=38, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_HEADING)
    box(sl, rx+0.1, 2.2, 2.8, 0.42, label, size=12,
        color=acc, align=PP_ALIGN.CENTER, font=FONT_BODY)

# Left — runs table (simulated)
rect(sl, 0.55, 2.85, 7.0, 4.05, RGBColor(0x04, 0x12, 0x28),
     border_color=RGBColor(0x2a, 0x4a, 0x80))
box(sl, 0.72, 2.92, 6.6, 0.36, "WORKFLOW RUNS  (/workflow-runs page)",
    size=12, bold=True, color=C_BLUE_ACC, font=FONT_HEADING)

# search bar mock
rect(sl, 0.72, 3.35, 4.5, 0.38, RGBColor(0x08, 0x1c, 0x40),
     border_color=RGBColor(0x2a, 0x4a, 0x80))
box(sl, 0.88, 3.39, 4.2, 0.3, "Search by input / output / run ID...",
    size=11, color=RGBColor(0x44, 0x66, 0x99), font=FONT_BODY, italic=True)
for label, bg, x in [("All", RGBColor(0x16, 0x52, 0x9a), 5.3),
                      ("Completed", RGBColor(0x08, 0x38, 0x18), 5.92),
                      ("Failed", RGBColor(0x38, 0x08, 0x08), 6.82)]:
    r2 = sl.shapes.add_shape(1, Inches(x), Inches(3.35), Inches(0.76), Inches(0.38))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg; r2.line.fill.background()
    box(sl, x+0.06, 3.39, 0.65, 0.3, label, size=10, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_BODY)

run_rows = [
    ("run_abc123", "Customer FAQ workflow", "2.1s", "Completed", C_GREEN_ACC),
    ("run_def456", "Document summarizer",   "3.8s", "Completed", C_GREEN_ACC),
    ("run_ghi789", "Lead qualifier agent",  "1.2s", "Failed",    C_RED_ACC),
    ("run_jkl012", "HR onboarding flow",    "4.5s", "Completed", C_GREEN_ACC),
    ("run_mno345", "Code review agent",     "2.9s", "Completed", C_GREEN_ACC),
]
for i, (rid, wf, dur, status, sc) in enumerate(run_rows):
    ry = 3.84 + i*0.56
    rbg = RGBColor(0x06, 0x16, 0x32) if i % 2 == 0 else RGBColor(0x08, 0x1c, 0x3e)
    rect(sl, 0.72, ry, 6.68, 0.5, rbg)
    box(sl, 0.82, ry+0.1, 1.6, 0.3, rid, size=11, color=C_BLUE_ACC, font=FONT_CODE)
    box(sl, 2.5,  ry+0.1, 2.8, 0.3, wf,  size=11, color=C_OFF_WHITE, font=FONT_BODY)
    box(sl, 5.35, ry+0.1, 0.7, 0.3, dur, size=11, color=C_MUTED_TXT, font=FONT_BODY, align=PP_ALIGN.RIGHT)
    box(sl, 6.1,  ry+0.1, 1.2, 0.3, status, size=11, bold=True, color=sc, font=FONT_BODY, align=PP_ALIGN.CENTER)

# Right — Trace drawer
rect(sl, 7.72, 2.85, 5.1, 4.05, RGBColor(0x08, 0x1a, 0x3a),
     border_color=C_BLUE_ACC)
box(sl, 7.88, 2.92, 4.7, 0.36, "EXECUTION TRACE DRAWER  (480px overlay)",
    size=11, bold=True, color=C_BLUE_ACC, font=FONT_HEADING)
box(sl, 7.88, 3.32, 4.7, 0.28, "Run: run_abc123  ·  Workflow: Customer FAQ",
    size=11, color=C_OFF_WHITE, font=FONT_BODY)
trace_nodes = [
    ("Input Received",    "0ms",   "done",    C_GREEN_ACC),
    ("RAG Query",         "210ms", "done",    C_GREEN_ACC),
    ("GPT-4o LLM Call",   "1.6s",  "done",    C_GREEN_ACC),
    ("Guardrails Check",  "88ms",  "done",    C_GREEN_ACC),
    ("Output Delivered",  "12ms",  "done",    C_GREEN_ACC),
]
for i, (node, dur, status, sc) in enumerate(trace_nodes):
    ny = 3.7 + i*0.62
    rect(sl, 7.88, ny, 4.76, 0.52, RGBColor(0x06, 0x14, 0x30),
         border_color=sc, border_w=0.5)
    box(sl, 8.0,  ny+0.06, 2.6, 0.28, node, size=12, bold=True, color=C_WHITE, font=FONT_BODY)
    box(sl, 10.7, ny+0.06, 0.8, 0.28, dur,  size=11, color=C_MUTED_TXT, font=FONT_BODY, align=PP_ALIGN.RIGHT)
    box(sl, 8.0,  ny+0.3,  4.5, 0.18, f"Status: {status}",
        size=10, color=sc, font=FONT_BODY)

box(sl, 0.55, 7.05, 12.23, 0.38,
    "Backed by workflow_runs PostgreSQL table  ·  Full JSONB node logs  ·  Fixed overlay drawer slides in from right at 480px width",
    size=12, color=RGBColor(0x44, 0x77, 0xcc), font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — vs Lyzr Architect (Detailed)
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_WHITE_BG)
box(sl, 0.55, 0.28, 5, 0.42, "COMPETITIVE COMPARISON", size=12, bold=True,
    color=C_MUTED_TXT, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "AgentForge vs Lyzr Architect — Feature Parity",
    size=28, bold=True, color=C_DARK_TXT, font=FONT_HEADING)

headers = ["Capability", "AgentForge  v4.0  ✅", "Lyzr Architect"]
rows_lyzr = [
    ["Visual drag-drop canvas",         "✅ ReactFlow v12",           "✅"],
    ["Prompt-to-agent (NL)",            "✅ GPT-4o powered",          "✅"],
    ["Multi-agent orchestration",       "✅ Manager/worker",          "✅"],
    ["Agent type filters (Studio)",     "✅ Agent/Mgr/Superflow",     "✅"],
    ["RAG pipeline",                    "✅ LangChain+pgvector",      "✅"],
    ["PII redaction",                   "✅ Microsoft Presidio",      "✅"],
    ["Safety rules dashboard",          "✅ 8 rules + toggle UI",     "Partial"],
    ["Evaluations / test engine",       "✅ Score + history",         "❌"],
    ["Voice agents (TTS/STT)",          "✅ Azure Speech SDK",        "❌"],
    ["Planning Architect",              "✅ 60s full-stack gen",      "❌"],
    ["Full-Stack ZIP Download",         "✅ React+FastAPI+Docker",    "❌"],
    ["Workflow Observability",          "✅ Trace drawer + stats",    "❌"],
    ["Azure-native deployment",         "✅ GPT-4o + Speech",         "❌ OpenAI SaaS only"],
    ["Self-hosted / open source",       "✅ Docker Compose",          "❌ SaaS only"],
    ["Agent versioning snapshots",      "✅ Immutable history",       "Partial"],
]
tbl(sl, 0.55, 1.58, 12.23, 5.58, headers, rows_lyzr, col_widths=[5.0, 3.8, 3.4])

# bottom note
r_note = sl.shapes.add_shape(1, Inches(0.55), Inches(7.08), Inches(12.23), Inches(0.36))
r_note.fill.solid(); r_note.fill.fore_color.rgb = RGBColor(0xe0, 0xec, 0xff)
r_note.line.color.rgb = RGBColor(0x16, 0x52, 0x9a); r_note.line.width = Pt(0.75)
box(sl, 0.72, 7.1, 12, 0.3,
    "AgentForge v4.0 adds Planning Architect, Full-Stack ZIP, Workflow Observability — not available in Lyzr.",
    size=13, bold=True, color=RGBColor(0x16, 0x52, 0x9a), font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Market Comparison (Broad)
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_SLATE_BG)
box(sl, 0.55, 0.28, 5, 0.42, "COMPETITIVE LANDSCAPE", size=12, bold=True,
    color=C_MUTED_TXT, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "AgentForge vs The Broader Market",
    size=28, bold=True, color=C_DARK_TXT, font=FONT_HEADING)

headers2 = ["Tool", "Canvas", "Guardrails", "Evaluations", "Voice", "Architect", "Azure-Native", "Self-Host", "Simulation"]
rows_market = [
    ["AgentForge",        "✅", "✅", "✅", "✅", "✅", "✅", "✅", "✅"],
    ["Lyzr Architect",    "✅", "Partial", "❌", "❌", "❌", "❌", "❌", "❌"],
    ["Flowise",           "✅", "❌", "❌", "❌", "❌", "❌", "✅", "❌"],
    ["Langflow",          "✅", "❌", "❌", "❌", "❌", "❌", "✅", "❌"],
    ["Dify",              "✅", "❌", "❌", "❌", "❌", "❌", "✅", "❌"],
    ["CrewAI",            "❌", "❌", "❌", "❌", "❌", "❌", "✅", "❌"],
    ["MS Copilot Studio", "✅", "✅", "❌", "❌", "❌", "✅", "❌", "❌"],
    ["n8n",               "✅", "❌", "❌", "❌", "❌", "❌", "✅", "❌"],
]
tbl(sl, 0.55, 1.58, 12.23, 5.3, headers2, rows_market,
    col_widths=[2.2, 1.06, 1.22, 1.22, 1.0, 1.18, 1.22, 1.06, 1.07])

r_note = sl.shapes.add_shape(1, Inches(0.55), Inches(7.05), Inches(12.23), Inches(0.4))
r_note.fill.solid(); r_note.fill.fore_color.rgb = RGBColor(0xe0, 0xec, 0xff)
r_note.line.color.rgb = RGBColor(0x16, 0x52, 0x9a); r_note.line.width = Pt(0.75)
box(sl, 0.72, 7.08, 12, 0.32,
    "AgentForge is the only platform with all 8 dimensions — including Planning Architect not available anywhere else.",
    size=13, bold=True, color=RGBColor(0x16, 0x52, 0x9a), font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Complete Feature List (20 features)
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_DARK_BG)
top_bar(sl, RGBColor(0x16, 0x52, 0x9a))
box(sl, 0.55, 0.28, 5, 0.42, "COMPLETE FEATURE SET", size=12, bold=True,
    color=C_BLUE_ACC, font=FONT_HEADING)
box(sl, 0.55, 0.68, 12, 0.78, "20 Enterprise Features Shipped",
    size=30, bold=True, color=C_WHITE, font=FONT_HEADING)

all_feats = [
    ("Visual Drag-Drop Builder",   "ReactFlow canvas, animated edges, dark mode",            C_BLUE_ACC),
    ("Prompt-to-Agent AI",         "GPT-4o generates full config from one sentence",          C_BLUE_ACC),
    ("Multi-Agent Orchestration",  "Manager delegates to workers, sequential/parallel",       C_GREEN_ACC),
    ("Agent Type Filters",         "Agent / Managerial / Superflow tabs in Studio",           C_BLUE_ACC),
    ("RAG Knowledge Pipeline",     "LangChain + pgvector, PDF/TXT ingest, citations",         C_GREEN_ACC),
    ("PII Redaction",              "Microsoft Presidio, 5 entity types, always-on",           C_RED_ACC),
    ("Safety Rules Dashboard",     "8 rules, live toggle, category/severity badges",          C_RED_ACC),
    ("Evaluations Engine",         "JSON test cases, score, history, pass-rate gate",         C_AMBER_ACC),
    ("Voice Agents",               "Azure Speech SDK, 8 Neural voices, TTS/STT config",       C_PURPLE_ACC),
    ("Planning Architect",         "NL -> full plan + React UI preview + deployable ZIP",     C_PURPLE_ACC),
    ("Full-Stack ZIP Download",    "React+FastAPI+FAISS+PostgreSQL+Docker in one export",     C_PURPLE_ACC),
    ("Workflow Observability",     "Trace drawer, stats bar, per-node logs, JSONB storage",   C_BLUE_ACC),
    ("Simulation Engine",          "Batch test runner, expected output scoring",              C_AMBER_ACC),
    ("Tool Registry",              "7 tools: web, calc, email, Slack, GitHub, Jira, Drive",   C_TEAL_ACC),
    ("Control Plane Dashboard",    "Live stats, guardrail triggers, avg latency",             C_BLUE_ACC),
    ("Agent Versioning",           "Immutable snapshots on every save, rollback ready",       C_BLUE_ACC),
    ("Audit Logs",                 "Full run history: inputs, outputs, guardrail events",     C_GREEN_ACC),
    ("RBAC + JWT Auth",            "Admin / Developer / Viewer, enforced on all routes",      C_RED_ACC),
    ("OpenTelemetry Tracing",      "Jaeger, Azure Monitor, GCP, Datadog, console/none",       C_TEAL_ACC),
    ("Docker Compose Stack",       "postgres+pgvector + backend + frontend + jaeger",         C_GREEN_ACC),
]
for i, (name, desc, acc) in enumerate(all_feats):
    col = i % 4; row = i // 4
    cx = 0.55 + col*3.2; cy = 1.72 + row*1.08
    r2 = sl.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(3.1), Inches(0.98))
    r2.fill.solid(); r2.fill.fore_color.rgb = RGBColor(0x0c, 0x1c, 0x3a)
    r2.line.color.rgb = acc; r2.line.width = Pt(0.5)
    a = sl.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(0.05), Inches(0.98))
    a.fill.solid(); a.fill.fore_color.rgb = acc; a.line.fill.background()
    box(sl, cx+0.12, cy+0.06, 2.86, 0.34, name, size=12, bold=True,
        color=C_WHITE, font=FONT_HEADING)
    box(sl, cx+0.12, cy+0.44, 2.86, 0.46, desc, size=11, color=C_OFF_WHITE, font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — Quality & Delivery Metrics
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_BLUE_BG)
top_bar(sl, C_BLUE_ACC)
slide_header(sl, "Delivery Quality", "Built to Production Standard")

metrics = [
    ("32", "Tests Passing\n(32/32)",       RGBColor(0x04, 0x18, 0x34)),
    ("20", "Core Features\nShipped",        RGBColor(0x04, 0x18, 0x34)),
    ("11", "API Route\nModules",            RGBColor(0x04, 0x18, 0x34)),
    ("3",  "RBAC Roles\nEnforced",          RGBColor(0x04, 0x18, 0x34)),
]
for i, (num, label, bg) in enumerate(metrics):
    rx = 0.55 + i*3.2
    r2 = sl.shapes.add_shape(1, Inches(rx), Inches(1.92), Inches(3.0), Inches(1.45))
    r2.fill.solid(); r2.fill.fore_color.rgb = bg
    r2.line.color.rgb = C_BLUE_ACC; r2.line.width = Pt(1.0)
    box(sl, rx+0.1, 1.97, 2.8, 0.78, num, size=48, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_HEADING)
    box(sl, rx+0.1, 2.72, 2.8, 0.58, label, size=13, color=C_OFF_WHITE,
        align=PP_ALIGN.CENTER, font=FONT_BODY)

# API endpoints
rect(sl, 0.55, 3.55, 6.1, 3.65, RGBColor(0x02, 0x0e, 0x22),
     border_color=RGBColor(0x1e, 0x40, 0x80))
box(sl, 0.72, 3.65, 5.8, 0.36, "KEY API ENDPOINTS",
    size=12, bold=True, color=RGBColor(0x55, 0x99, 0xdd), font=FONT_HEADING)
endpoints = [
    "POST /api/agents/generate          ->  NL-to-agent config",
    "POST /api/agents/{id}/run          ->  Execute with guardrails",
    "POST /api/rag/ingest               ->  Upload & chunk documents",
    "POST /api/simulation/run           ->  Batch test runner",
    "GET  /api/safety/rules             ->  List / toggle safety rules",
    "POST /api/evaluations/runs         ->  Run evaluation suite",
    "POST /api/voice/synthesize         ->  Azure TTS audio stream",
    "POST /api/architect/generate       ->  Planning Architect plan",
    "GET  /api/workflow-runs            ->  All workflow run traces",
]
for i, ep in enumerate(endpoints):
    box(sl, 0.72, 4.08 + i*0.38, 5.8, 0.34, ep, size=11, color=C_OFF_WHITE,
        font=FONT_CODE)

# Security items
rect(sl, 6.85, 3.55, 6.0, 3.65, RGBColor(0x02, 0x0e, 0x22),
     border_color=RGBColor(0x1e, 0x40, 0x80))
box(sl, 7.0, 3.65, 5.7, 0.36, "SECURITY & QUALITY",
    size=12, bold=True, color=RGBColor(0x55, 0x99, 0xdd), font=FONT_HEADING)
sec_items = [
    "No hardcoded secrets — config fails loud at startup",
    "Azure API keys git-ignored via .env + start.ps1",
    "bcrypt==4.0.1 pinned for passlib compatibility",
    "PostgreSQL healthcheck before backend starts",
    "SQLAlchemy 2.0 async end-to-end",
    "Pydantic v2 validation at every API boundary",
    "JWT enforced on all 11 router modules",
    "FAISS index rebuilt on document upload (Planning Architect)",
]
for i, s in enumerate(sec_items):
    box(sl, 7.0, 4.1 + i*0.44, 5.7, 0.38, f"✓  {s}", size=12,
        color=C_GREEN_ACC, font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 18 — Use Cases
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_AMBER_BG)
top_bar(sl, C_AMBER_ACC, 0.08)
slide_header(sl, "Real-World Use Cases", "Where AgentForge Shines", title_color=C_WHITE)

cases = [
    ("Financial Document Q&A",
     "Ingest annual reports via RAG. Agent answers investor questions grounded in source data. "
     "PII guardrails block SSN/account number disclosure. Audit log for every query.",
     C_AMBER_ACC),
    ("Healthcare Triage Assistant",
     "Intake agent -> diagnosis agent -> scheduling agent pipeline. "
     "Presidio removes patient PII before any LLM call. HIPAA-compliant by design.",
     C_AMBER_ACC),
    ("DevOps Automation",
     "Manager delegates: GitHub reads PRs, Jira creates tickets, Slack posts summaries — "
     "one multi-agent workflow replaces three separate automation scripts.",
     C_GREEN_ACC),
    ("Rapid App Prototyping",
     "Use Planning Architect to describe your app — get a full project plan, working React UI preview, "
     "and a deployable ZIP in 60 seconds. From idea to running app without an engineering team.",
     C_PURPLE_ACC),
    ("Enterprise Knowledge Base",
     "HR, legal, and compliance docs ingested via RAG. Employees ask questions in plain English; "
     "agents answer with source citations. Safety rules prevent off-topic or harmful responses.",
     C_TEAL_ACC),
    ("Pre-Production QA Gate",
     "Evaluation engine scores agents on a regression test suite overnight. "
     "Only agents meeting 90%+ pass rate are promoted to production — zero manual QA effort.",
     C_BLUE_ACC),
]
for i, (title, body, acc) in enumerate(cases):
    col = i % 3; row = i // 2
    card_dark(sl, 0.55 + col*4.27, 2.0 + row*2.52, 4.0, 2.38,
              title, body, accent=acc, title_size=15, body_size=12)


# ════════════════════════════════════════════════════════════════
# SLIDE 19 — Roadmap
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_DARK_BG)
top_bar(sl, RGBColor(0x16, 0x52, 0x9a))
slide_header(sl, "Strategic Roadmap", "What Comes Next — AgentForge")

roadmap = [
    ("1", "Azure AI Search Vector RAG",
     "Replace in-memory chunks with Azure AI Search for enterprise-scale semantic retrieval "
     "with full vector similarity, faceting, and hybrid search.",
     RGBColor(0x16, 0x52, 0x9a), C_BLUE_ACC),
    ("2", "WebSocket Streaming Responses",
     "Real-time token-by-token agent output. Azure client streaming is wired — "
     "frontend upgrade only. Sub-second first-token latency.",
     RGBColor(0x4a, 0x30, 0xb0), C_PURPLE_ACC),
    ("3", "Live Tool Credentials UI",
     "Connect real Slack, GitHub, and email accounts via a settings panel. "
     "Move integrations from stubs to fully live without code changes.",
     RGBColor(0x0c, 0x5e, 0x48), C_TEAL_ACC),
    ("4", "Telephony Inbound Calls",
     "Wire Azure Communication Services to Voice Agents for real inbound call handling. "
     "IVR flows, call recording, and transcript storage.",
     RGBColor(0x4a, 0x30, 0xb0), C_PURPLE_ACC),
    ("5", "Multi-Tenant Workspaces",
     "Isolated workspaces per team — separate agents, RAG indexes, audit logs, "
     "safety configs, and quotas per tenant.",
     RGBColor(0x7a, 0x20, 0x20), C_RED_ACC),
    ("6", "In-Browser Agent Chat",
     "Test any agent in a chat window directly in the builder — "
     "no API client needed. Instant feedback loop for prompt engineers.",
     RGBColor(0x78, 0x46, 0x08), C_AMBER_ACC),
]
for i, (num, title, body, bg, acc) in enumerate(roadmap):
    col = i % 3; row = i // 3
    cx = 0.55 + col*4.27; cy = 2.0 + row*2.6
    card_dark(sl, cx, cy, 4.0, 2.45, f"{num}.  {title}", body,
              accent=acc, title_size=15, body_size=12)

# Shipped banner
rect(sl, 0.55, 7.0, 12.23, 0.44, RGBColor(0x04, 0x28, 0x14),
     border_color=C_GREEN_ACC)
box(sl, 0.72, 7.04, 12, 0.34,
    "SHIPPED IN v4.0:  Planning Architect  ·  Full-Stack ZIP Download  ·  Workflow Observability  ·  Voice Agent Azure TTS  ·  Evaluations Engine",
    size=12, bold=True, color=C_GREEN_ACC, font=FONT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 20 — Get Started / Closing
# ════════════════════════════════════════════════════════════════
sl = add_slide(C_TEAL_BG)
top_bar(sl, RGBColor(0x18, 0x7a, 0x5e), 0.08)
box(sl, 0.65, 0.45, 5, 0.45, "DEPLOY TODAY", size=13, bold=True,
    color=C_TEAL_ACC, font=FONT_HEADING)
box(sl, 0.65, 0.88, 10, 1.1, "Up and Running in 3 Commands",
    size=40, bold=True, color=C_WHITE, font=FONT_HEADING)

steps_cmds = [
    ("1 — CONFIGURE",  "cp .env.example .env        # add your Azure OpenAI + Speech keys"),
    ("2 — LAUNCH",     "docker-compose up --build"),
    ("3 — OPEN",       "http://localhost:5173              # frontend   |   :8000/docs  ->  API explorer"),
]
for i, (label, cmd) in enumerate(steps_cmds):
    rect(sl, 0.65, 2.2 + i*1.32, 11.5, 1.18,
         RGBColor(0x01, 0x16, 0x0e), border_color=RGBColor(0x18, 0x7a, 0x5e))
    box(sl, 0.85, 2.24 + i*1.32, 4, 0.38, label,
        size=12, bold=True, color=C_TEAL_ACC, font=FONT_HEADING)
    box(sl, 0.85, 2.6 + i*1.32, 11.1, 0.55, cmd,
        size=15, bold=True, color=C_WHITE, font=FONT_CODE)

divider(sl, 6.2, color=RGBColor(0x18, 0x7a, 0x5e))

info_pills = [
    ("Project:  C:\\Users\\n.sureshmanikandan\\Repo1\\AgentForge", RGBColor(0x01, 0x1e, 0x14), C_TEAL_ACC),
    ("Version:  4.0  —  July 2026",                               RGBColor(0x01, 0x1e, 0x14), C_TEAL_ACC),
    ("Features:  20  ·  Tests:  32/32",                           RGBColor(0x01, 0x1e, 0x14), C_GREEN_ACC),
]
for i, (txt, bg, fg) in enumerate(info_pills):
    pill(sl, 0.65 + i*4.3, 6.38, 4.0, 0.5, txt, bg, fg, 12)

# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
out = r"C:\Users\n.sureshmanikandan\Repo1\AgentForge\docs\AgentForge-PresentationV4.0.pptx"
prs.save(out)
print(f"\nSaved: {out}")
print(f"    Slides: {len(prs.slides)}")
